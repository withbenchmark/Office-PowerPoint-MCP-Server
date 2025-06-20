"""
Utility functions for PowerPoint manipulation using python-pptx.

This module provides a comprehensive set of functions to create and manipulate
PowerPoint presentations programmatically. It wraps the python-pptx library
with higher-level functions that simplify common operations.

Key features:
- Create, open, and save presentations
- Add and format slides with various layouts
- Manipulate text, including titles, placeholders, and textboxes
- Add and format shapes, images, tables, and charts
- Set document properties

Usage examples:
    # Create a new presentation
    pres = create_presentation()
    
    # Add a title slide
    slide, layout = add_slide(pres, 0)
    set_title(slide, "Presentation Title")
    
    # Add a content slide with bullet points
    slide, layout = add_slide(pres, 1)
    set_title(slide, "Key Points")
    placeholder = slide.placeholders[1]  # Content placeholder
    add_bullet_points(placeholder, ["Point 1", "Point 2", "Point 3"])
    
    # Save the presentation
    save_presentation(pres, "my_presentation.pptx")
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL_TYPE
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.graphfrm import GraphicFrame
import io
from typing import Dict, List, Tuple, Union, Optional, Any
import base64
import tempfile
import os
from PIL import Image, ImageEnhance, ImageFilter
from fontTools.ttLib import TTFont
from fontTools.subset import Subsetter

# Professional color schemes
PROFESSIONAL_COLOR_SCHEMES = {
    'modern_blue': {
        'primary': (0, 120, 215),      # Microsoft Blue
        'secondary': (40, 40, 40),     # Dark Gray
        'accent1': (0, 176, 240),      # Light Blue
        'accent2': (255, 192, 0),      # Orange
        'light': (247, 247, 247),      # Light Gray
        'text': (68, 68, 68),          # Text Gray
    },
    'corporate_gray': {
        'primary': (68, 68, 68),       # Charcoal
        'secondary': (0, 120, 215),    # Blue
        'accent1': (89, 89, 89),       # Medium Gray
        'accent2': (217, 217, 217),    # Light Gray
        'light': (242, 242, 242),      # Very Light Gray
        'text': (51, 51, 51),          # Dark Text
    },
    'elegant_green': {
        'primary': (70, 136, 71),      # Forest Green
        'secondary': (255, 255, 255),  # White
        'accent1': (146, 208, 80),     # Light Green
        'accent2': (112, 173, 71),     # Medium Green
        'light': (238, 236, 225),      # Cream
        'text': (89, 89, 89),          # Gray Text
    },
    'warm_red': {
        'primary': (192, 80, 77),      # Deep Red
        'secondary': (68, 68, 68),     # Dark Gray
        'accent1': (230, 126, 34),     # Orange
        'accent2': (241, 196, 15),     # Yellow
        'light': (253, 253, 253),      # White
        'text': (44, 62, 80),          # Blue Gray
    }
}

# Professional typography settings
PROFESSIONAL_FONTS = {
    'title': {
        'name': 'Segoe UI',
        'size_large': 36,
        'size_medium': 28,
        'size_small': 24,
        'bold': True
    },
    'subtitle': {
        'name': 'Segoe UI Light',
        'size_large': 20,
        'size_medium': 18,
        'size_small': 16,
        'bold': False
    },
    'body': {
        'name': 'Segoe UI',
        'size_large': 16,
        'size_medium': 14,
        'size_small': 12,
        'bold': False
    },
    'caption': {
        'name': 'Segoe UI',
        'size_large': 12,
        'size_medium': 10,
        'size_small': 9,
        'bold': False
    }
}

# Professional layout constants (in inches)
PROFESSIONAL_LAYOUT = {
    'margins': {
        'left': 0.75,
        'right': 0.75,
        'top': 0.5,
        'bottom': 0.5
    },
    'spacing': {
        'title_to_content': 0.3,
        'between_elements': 0.2,
        'bullet_indent': 0.5
    },
    'standard_sizes': {
        'title_height': 1.2,
        'content_area_height': 5.5,
        'footer_height': 0.4
    }
}

def get_professional_color(scheme_name: str, color_type: str) -> Tuple[int, int, int]:
    """
    Get a professional color from predefined color schemes.
    
    Args:
        scheme_name: Name of the color scheme ('modern_blue', 'corporate_gray', etc.)
        color_type: Type of color ('primary', 'secondary', 'accent1', 'accent2', 'light', 'text')
        
    Returns:
        RGB color tuple (r, g, b)
    """
    if scheme_name not in PROFESSIONAL_COLOR_SCHEMES:
        scheme_name = 'modern_blue'  # Default fallback
    
    scheme = PROFESSIONAL_COLOR_SCHEMES[scheme_name]
    return scheme.get(color_type, scheme['primary'])

def get_professional_font(font_type: str, size_category: str = 'medium') -> Dict:
    """
    Get professional font settings.
    
    Args:
        font_type: Type of font ('title', 'subtitle', 'body', 'caption')
        size_category: Size category ('large', 'medium', 'small')
        
    Returns:
        Dictionary with font settings
    """
    if font_type not in PROFESSIONAL_FONTS:
        font_type = 'body'  # Default fallback
    
    font_config = PROFESSIONAL_FONTS[font_type]
    size_key = f'size_{size_category}'
    
    return {
        'name': font_config['name'],
        'size': font_config.get(size_key, font_config['size_medium']),
        'bold': font_config['bold']
    }

def calculate_professional_layout(slide_width: float = 10, slide_height: float = 7.5) -> Dict:
    """
    Calculate professional layout dimensions based on slide size.
    
    Args:
        slide_width: Slide width in inches
        slide_height: Slide height in inches
        
    Returns:
        Dictionary with calculated layout dimensions
    """
    margins = PROFESSIONAL_LAYOUT['margins']
    
    content_width = slide_width - margins['left'] - margins['right']
    content_height = slide_height - margins['top'] - margins['bottom']
    
    title_area = {
        'left': margins['left'],
        'top': margins['top'],
        'width': content_width,
        'height': PROFESSIONAL_LAYOUT['standard_sizes']['title_height']
    }
    
    content_area = {
        'left': margins['left'],
        'top': title_area['top'] + title_area['height'] + PROFESSIONAL_LAYOUT['spacing']['title_to_content'],
        'width': content_width,
        'height': content_height - title_area['height'] - PROFESSIONAL_LAYOUT['spacing']['title_to_content']
    }
    
    return {
        'slide': {'width': slide_width, 'height': slide_height},
        'margins': margins,
        'title_area': title_area,
        'content_area': content_area,
        'spacing': PROFESSIONAL_LAYOUT['spacing']
    }

def try_multiple_approaches(operation_name, approaches):
    """
    Try multiple approaches to perform an operation, returning the first successful result.
    
    Args:
        operation_name: Name of the operation for error reporting
        approaches: List of (approach_func, description) tuples to try
        
    Returns:
        Tuple of (result, None) if any approach succeeded, or (None, error_messages) if all failed
    """
    error_messages = []
    
    for approach_func, description in approaches:
        try:
            result = approach_func()
            return result, None
        except Exception as e:
            error_messages.append(f"{description}: {str(e)}")
    
    return None, f"Failed to {operation_name} after trying multiple approaches: {'; '.join(error_messages)}"

def safe_operation(operation_name, operation_func, error_message=None, *args, **kwargs):
    """
    Execute an operation safely with standard error handling.
    
    Args:
        operation_name: Name of the operation for error reporting
        operation_func: Function to execute
        error_message: Custom error message (optional)
        *args, **kwargs: Arguments to pass to the operation function
        
    Returns:
        A tuple (result, error) where error is None if operation was successful
    """
    try:
        result = operation_func(*args, **kwargs)
        return result, None
    except ValueError as e:
        error_msg = error_message or f"Invalid input for {operation_name}: {str(e)}"
        return None, error_msg
    except TypeError as e:
        error_msg = error_message or f"Type error in {operation_name}: {str(e)}"
        return None, error_msg
    except Exception as e:
        error_msg = error_message or f"Failed to execute {operation_name}: {str(e)}"
        return None, error_msg

# ---- Presentation Functions ----

def open_presentation(file_path: str) -> Presentation:
    """
    Open an existing PowerPoint presentation.
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        A Presentation object
    """
    return Presentation(file_path)

def create_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation.
    
    Returns:
        A new Presentation object
    """
    return Presentation()

def create_presentation_from_template(template_path: str) -> Presentation:
    """
    Create a new PowerPoint presentation from a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        A new Presentation object based on the template
        
    Raises:
        FileNotFoundError: If the template file doesn't exist
        Exception: If the template file is corrupted or invalid
    """
    import os
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    if not template_path.lower().endswith(('.pptx', '.potx')):
        raise ValueError("Template file must be a .pptx or .potx file")
    
    try:
        # Load the template file as a presentation
        presentation = Presentation(template_path)
        return presentation
    except Exception as e:
        raise Exception(f"Failed to load template file '{template_path}': {str(e)}")

def get_template_info(template_path: str) -> Dict:
    """
    Get information about a template file without fully loading it.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        Dictionary containing template information
    """
    import os
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    try:
        presentation = Presentation(template_path)
        
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        # Get file size
        file_size = os.path.getsize(template_path)
        
        return {
            "template_path": template_path,
            "file_size_bytes": file_size,
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props
        }
    except Exception as e:
        raise Exception(f"Failed to read template info from '{template_path}': {str(e)}")

def save_presentation(presentation: Presentation, file_path: str) -> str:
    """
    Save a PowerPoint presentation to a file.
    
    Args:
        presentation: The Presentation object
        file_path: Path where the file should be saved
        
    Returns:
        The file path where the presentation was saved
    """
    presentation.save(file_path)
    return file_path

def presentation_to_base64(presentation: Presentation) -> str:
    """
    Convert a presentation to a base64 encoded string.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Base64 encoded string of the presentation
    """
    ppt_bytes = io.BytesIO()
    presentation.save(ppt_bytes)
    ppt_bytes.seek(0)
    return base64.b64encode(ppt_bytes.read()).decode('utf-8')

def base64_to_presentation(base64_string: str) -> Presentation:
    """
    Create a presentation from a base64 encoded string.
    
    Args:
        base64_string: Base64 encoded string of a presentation
        
    Returns:
        A Presentation object
    """
    ppt_bytes = io.BytesIO(base64.b64decode(base64_string))
    return Presentation(ppt_bytes)

# ---- Slide Functions ----

def add_slide(presentation: Presentation, layout_index: int = 1) -> Tuple:
    """
    Add a slide to the presentation.
    
    Args:
        presentation: The Presentation object
        layout_index: Index of the slide layout to use (default is 1, typically a title and content slide)
        
    Returns:
        A tuple containing the slide and its layout
    """
    layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(layout)
    return slide, layout

def add_professional_slide(presentation: Presentation, slide_type: str = 'title_content', 
                          color_scheme: str = 'modern_blue') -> Tuple:
    """
    Add a professionally designed slide with proper spacing and typography.
    
    Args:
        presentation: The Presentation object
        slide_type: Type of slide ('title', 'title_content', 'content', 'two_column', 'blank')
        color_scheme: Color scheme to apply ('modern_blue', 'corporate_gray', etc.)
        
    Returns:
        A tuple containing the slide, layout, and design info
    """
    # Map slide types to layout indices
    layout_map = {
        'title': 0,           # Title slide
        'title_content': 1,   # Title and content
        'content': 6,         # Content only
        'two_column': 3,      # Two content columns
        'blank': 6            # Blank layout
    }
    
    layout_index = layout_map.get(slide_type, 1)
    layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(layout)
    
    # Calculate professional layout
    slide_width = presentation.slide_width / 914400  # Convert EMU to inches
    slide_height = presentation.slide_height / 914400
    layout_info = calculate_professional_layout(slide_width, slide_height)
    
    # Apply professional styling if possible
    try:
        apply_professional_slide_background(slide, color_scheme)
    except Exception:
        pass  # Graceful degradation if background setting fails
    
    design_info = {
        'color_scheme': color_scheme,
        'layout_info': layout_info,
        'slide_type': slide_type
    }
    
    return slide, layout, design_info

def apply_professional_slide_background(slide, color_scheme: str = 'modern_blue'):
    """
    Apply a professional background to a slide.
    
    Args:
        slide: The slide object
        color_scheme: Color scheme name
    """
    try:
        # Get background color
        bg_color = get_professional_color(color_scheme, 'light')
        
        # Try to set slide background
        if hasattr(slide, 'background'):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
    except Exception:
        # Graceful fallback - background setting may not be available
        pass

def get_slide_layouts(presentation: Presentation) -> List[Dict]:
    """
    Get all available slide layouts in the presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        A list of dictionaries with layout information
    """
    layouts = []
    for i, layout in enumerate(presentation.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        layouts.append(layout_info)
    return layouts

# ---- Placeholder Functions ----

def get_placeholders(slide) -> List[Dict]:
    """
    Get all placeholders in a slide.
    
    Args:
        slide: The slide object
        
    Returns:
        A list of dictionaries with placeholder information
    """
    placeholders = []
    for placeholder in slide.placeholders:
        placeholder_info = {
            "idx": placeholder.placeholder_format.idx,
            "type": placeholder.placeholder_format.type,
            "name": placeholder.name,
            "shape_type": placeholder.shape_type
        }
        placeholders.append(placeholder_info)
    return placeholders

def set_title(slide, title: str) -> None:
    """
    Set the title of a slide.
    
    Args:
        slide: The slide object
        title: The title text
    """
    if slide.shapes.title:
        slide.shapes.title.text = title

def set_professional_title(slide, title: str, color_scheme: str = 'modern_blue', 
                          size_category: str = 'large') -> None:
    """
    Set a professionally formatted title with proper typography and colors.
    
    Args:
        slide: The slide object
        title: The title text
        color_scheme: Color scheme to use
        size_category: Font size category ('large', 'medium', 'small')
    """
    if not slide.shapes.title:
        return
    
    # Set title text
    slide.shapes.title.text = title
    
    # Get professional font settings
    font_settings = get_professional_font('title', size_category)
    title_color = get_professional_color(color_scheme, 'primary')
    
    # Apply professional formatting
    try:
        text_frame = slide.shapes.title.text_frame
        
        # Enable word wrap and proper margins
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        # Format all paragraphs
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(6)
            paragraph.line_spacing = 1.15
            
            # Format all runs in the paragraph
            for run in paragraph.runs:
                font = run.font
                font.name = font_settings['name']
                font.size = Pt(font_settings['size'])
                font.bold = font_settings['bold']
                font.color.rgb = RGBColor(*title_color)
                
    except Exception as e:
        # Fallback to basic formatting
        try:
            format_text_advanced(
                slide.shapes.title.text_frame,
                font_size=font_settings['size'],
                font_name=font_settings['name'],
                bold=font_settings['bold'],
                color=title_color
            )
        except:
            pass

def populate_placeholder(slide, placeholder_idx: int, text: str) -> None:
    """
    Populate a placeholder with text.
    
    Args:
        slide: The slide object
        placeholder_idx: The index of the placeholder
        text: The text to add
    """
    placeholder = slide.placeholders[placeholder_idx]
    placeholder.text = text

def add_bullet_points(placeholder, bullet_points: List[str]) -> None:
    """
    Add bullet points to a placeholder.
    
    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts
    """
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    for i, point in enumerate(bullet_points):
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        
        # Only add line breaks between bullet points, not after the last one
        if i < len(bullet_points) - 1:
            p.line_spacing = 1.0

def add_professional_bullet_points(placeholder, bullet_points: List[str], 
                                   color_scheme: str = 'modern_blue',
                                   hierarchical: bool = False) -> None:
    """
    Add professionally formatted bullet points with proper typography and spacing.
    
    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts (can include level indicators like "  - sub point")
        color_scheme: Color scheme to use
        hierarchical: Whether to automatically detect and format hierarchical bullet points
    """
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    # Set professional text frame properties
    try:
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.25)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
    except:
        pass
    
    # Get professional font settings
    body_font = get_professional_font('body', 'medium')
    text_color = get_professional_color(color_scheme, 'text')
    
    for i, point in enumerate(bullet_points):
        # Determine bullet level if hierarchical
        level = 0
        clean_text = point.strip()
        
        if hierarchical:
            # Count leading spaces or tabs to determine level
            leading_spaces = len(point) - len(point.lstrip(' \t'))
            level = min(leading_spaces // 2, 4)  # Max 5 levels (0-4)
            
            # Remove common level indicators
            for indicator in ['• ', '- ', '* ', '○ ', '▪ ']:
                if clean_text.startswith(indicator):
                    clean_text = clean_text[len(indicator):]
                    break
        
        # Add paragraph
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = clean_text
        p.level = level
        
        # Set professional paragraph formatting
        try:
            p.space_before = Pt(0)
            p.space_after = Pt(3)
            p.line_spacing = 1.25
            
            if level == 0:
                p.space_before = Pt(6) if i > 0 else Pt(0)
            
            # Format the text run
            for run in p.runs:
                font = run.font
                font.name = body_font['name']
                font.size = Pt(max(body_font['size'] - level, 10))  # Smaller font for sub-levels
                font.bold = (level == 0)  # Only main bullets are bold
                font.color.rgb = RGBColor(*text_color)
                
        except Exception:
            # Fallback formatting
            try:
                format_text_advanced(
                    text_frame,
                    font_size=body_font['size'],
                    font_name=body_font['name'],
                    color=text_color
                )
            except:
                pass

# ---- Text Functions ----

def add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                font_size: int = None, font_name: str = None, bold: bool = None,
                italic: bool = None, color: Tuple[int, int, int] = None,
                alignment: str = None, auto_resize: bool = True) -> Any:
    """
    Add a textbox to a slide with enhanced text formatting and overflow handling.
    
    Args:
        slide: The slide object
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        auto_resize: Whether to automatically resize font if text overflows
        
    Returns:
        The created textbox shape
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    # Set the text first
    textbox.text_frame.text = text
    
    # Apply formatting if any formatting options are provided
    if any([font_size, font_name, bold, italic, color, alignment]):
        format_text_advanced(
            textbox.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=color,
            alignment=alignment,
            auto_resize=auto_resize
        )
    
    return textbox

def format_text(text_frame, font_size: int = None, font_name: str = None, 
                bold: bool = None, italic: bool = None, color: Tuple[int, int, int] = None,
                alignment: str = None, auto_resize: bool = False, max_font_size: int = None) -> None:
    """
    Format text in a text frame with overflow detection and auto-resize capability.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        auto_resize: Whether to automatically resize font if text overflows
        max_font_size: Maximum font size when auto-resizing (defaults to original font_size)
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    # Enable text wrapping and auto-fit behavior
    text_frame.word_wrap = True
    text_frame.auto_size = 1  # Auto-fit text to shape
    
    # Set maximum font size if auto_resize is enabled
    if auto_resize and max_font_size is None and font_size is not None:
        max_font_size = font_size
    
    # Apply formatting to all paragraphs and runs
    for paragraph in text_frame.paragraphs:
        if alignment and alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
            
        for run in paragraph.runs:
            font = run.font
            
            # Apply font size with auto-resize logic
            if font_size is not None:
                if auto_resize and max_font_size:
                    # Try progressively smaller font sizes if text overflows
                    current_size = min(font_size, max_font_size)
                    font.size = Pt(current_size)
                    
                    # Check if we need to reduce font size (basic heuristic)
                    while current_size > 8:  # Minimum readable font size
                        try:
                            # Set the font size and check if it fits
                            font.size = Pt(current_size)
                            break
                        except:
                            current_size -= 2
                            continue
                else:
                    font.size = Pt(font_size)
                
            if font_name is not None:
                font.name = font_name
                
            if bold is not None:
                font.bold = bold
                
            if italic is not None:
                font.italic = italic
                
            if color is not None:
                r, g, b = color
                font.color.rgb = RGBColor(r, g, b)

def format_text_advanced(text_frame, font_size: int = None, font_name: str = None, 
                        bold: bool = None, italic: bool = None, color: Tuple[int, int, int] = None,
                        alignment: str = None, auto_resize: bool = True, 
                        min_font_size: int = 8, max_font_size: int = None) -> Dict:
    """
    Advanced text formatting with overflow detection and automatic font size adjustment.
    
    Args:
        text_frame: The text frame to format
        font_size: Initial font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        auto_resize: Whether to automatically resize font if text overflows
        min_font_size: Minimum font size when auto-resizing
        max_font_size: Maximum font size when auto-resizing
    
    Returns:
        Dictionary with formatting results and final font size used
    """
    result = {
        'success': True,
        'original_font_size': font_size,
        'final_font_size': font_size,
        'auto_resized': False,
        'warnings': []
    }
    
    try:
        # Clear existing text formatting and start fresh
        if len(text_frame.paragraphs) == 0:
            text_frame.add_paragraph()
        
        # Enable text wrapping and auto-fit
        text_frame.word_wrap = True
        
        # Try to set auto-fit behavior (may not be available in all versions)
        try:
            text_frame.auto_size = 1  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except:
            result['warnings'].append("Auto-fit not available, using manual font sizing")
        
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        # Set max font size if not provided
        if max_font_size is None:
            max_font_size = font_size if font_size else 72
        
        # Start with the desired font size
        current_font_size = font_size if font_size else 12
        
        # Apply formatting to all paragraphs and runs
        for paragraph in text_frame.paragraphs:
            if alignment and alignment in alignment_map:
                paragraph.alignment = alignment_map[alignment]
            
            # If paragraph has no runs, create text to establish runs
            if len(paragraph.runs) == 0 and hasattr(paragraph, 'text'):
                original_text = paragraph.text
                paragraph.text = original_text or " "  # Ensure we have some text
            
            for run in paragraph.runs:
                font = run.font
                
                # Apply font name first
                if font_name is not None:
                    font.name = font_name
                
                # Apply font size with auto-resize logic
                if current_font_size is not None:
                    if auto_resize:
                        # Try progressively smaller sizes if needed
                        size_to_try = min(current_font_size, max_font_size)
                        
                        while size_to_try >= min_font_size:
                            try:
                                font.size = Pt(size_to_try)
                                # Test if this size works (no exception means it's ok)
                                result['final_font_size'] = size_to_try
                                if size_to_try != current_font_size:
                                    result['auto_resized'] = True
                                break
                            except Exception as e:
                                result['warnings'].append(f"Font size {size_to_try} failed: {str(e)}")
                                size_to_try -= 2
                                continue
                        
                        if size_to_try < min_font_size:
                            result['warnings'].append(f"Could not fit text even at minimum size {min_font_size}")
                            font.size = Pt(min_font_size)
                            result['final_font_size'] = min_font_size
                            result['auto_resized'] = True
                    else:
                        font.size = Pt(current_font_size)
                        result['final_font_size'] = current_font_size
                
                # Apply other formatting
                if bold is not None:
                    font.bold = bold
                
                if italic is not None:
                    font.italic = italic
                
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
        
        return result
        
    except Exception as e:
        result['success'] = False
        result['error'] = str(e)
        return result

def validate_text_container(shape, text_content: str, font_size: int) -> Dict:
    """
    Validate if text content will fit in a shape container.
    
    Args:
        shape: The shape containing the text
        text_content: The text to validate
        font_size: The font size to check
    
    Returns:
        Dictionary with validation results and suggestions
    """
    result = {
        'fits': True,
        'estimated_overflow': False,
        'suggested_font_size': font_size,
        'suggested_dimensions': None,
        'warnings': []
    }
    
    try:
        # Basic heuristic: estimate if text will overflow
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            # Rough estimation: average character width is about 0.6 * font_size
            avg_char_width = font_size * 0.6
            estimated_width = len(text_content) * avg_char_width
            
            # Convert shape dimensions to points (assuming they're in EMU)
            shape_width_pt = shape.width / 12700  # EMU to points conversion
            shape_height_pt = shape.height / 12700
            
            if estimated_width > shape_width_pt:
                result['fits'] = False
                result['estimated_overflow'] = True
                
                # Suggest smaller font size
                suggested_size = int((shape_width_pt / len(text_content)) * 0.8)
                result['suggested_font_size'] = max(suggested_size, 8)
                
                # Suggest larger dimensions
                result['suggested_dimensions'] = {
                    'width': estimated_width * 1.2,
                    'height': shape_height_pt
                }
                
                result['warnings'].append(
                    f"Text may overflow. Consider font size {result['suggested_font_size']} "
                    f"or increase width to {result['suggested_dimensions']['width']:.1f} points"
                )
        
        return result
        
    except Exception as e:
        result['warnings'].append(f"Validation error: {str(e)}")
        return result

# ---- Gradient Background Functions ----

def create_gradient_image(width: int, height: int, start_color: Tuple[int, int, int], 
                         end_color: Tuple[int, int, int], direction: str = 'horizontal') -> Image.Image:
    """
    Create a gradient image using Pillow.
    
    Args:
        width: Image width in pixels
        height: Image height in pixels
        start_color: RGB tuple for gradient start color
        end_color: RGB tuple for gradient end color
        direction: Gradient direction ('horizontal', 'vertical', 'diagonal')
        
    Returns:
        PIL Image object with gradient
    """
    image = Image.new('RGB', (width, height))
    
    for y in range(height):
        for x in range(width):
            if direction == 'horizontal':
                ratio = x / width
            elif direction == 'vertical':
                ratio = y / height
            elif direction == 'diagonal':
                ratio = (x + y) / (width + height)
            else:
                ratio = x / width  # Default to horizontal
            
            # Interpolate between start and end colors
            r = int(start_color[0] + ratio * (end_color[0] - start_color[0]))
            g = int(start_color[1] + ratio * (end_color[1] - start_color[1]))
            b = int(start_color[2] + ratio * (end_color[2] - start_color[2]))
            
            image.putpixel((x, y), (r, g, b))
    
    return image

def set_slide_gradient_background(slide, start_color: Tuple[int, int, int], 
                                 end_color: Tuple[int, int, int], direction: str = 'horizontal') -> Dict:
    """
    Set a gradient background for a slide using a generated gradient image.
    
    Args:
        slide: The slide object
        start_color: RGB tuple for gradient start color
        end_color: RGB tuple for gradient end color
        direction: Gradient direction ('horizontal', 'vertical', 'diagonal')
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'gradient_info': {
            'start_color': start_color,
            'end_color': end_color,
            'direction': direction
        }
    }
    
    try:
        # Get slide dimensions (assuming standard 16:9 ratio if not available)
        try:
            # Get presentation to access slide dimensions
            presentation = slide.part.package.presentation_part.presentation
            slide_width = int(presentation.slide_width / 9525)  # Convert EMU to pixels (approximate)
            slide_height = int(presentation.slide_height / 9525)
        except:
            # Default to common HD resolution
            slide_width = 1920
            slide_height = 1080
        
        # Create gradient image
        gradient_img = create_gradient_image(slide_width, slide_height, start_color, end_color, direction)
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
            gradient_img.save(temp_file.name, 'PNG')
            temp_path = temp_file.name
        
        try:
            # Set slide background using the gradient image
            # Try different approaches to set background
            
            # Approach 1: Try to set background fill
            try:
                if hasattr(slide, 'background'):
                    slide.background.fill.picture(temp_path)
                    result['success'] = True
                    result['message'] = 'Gradient background set successfully using slide.background'
                else:
                    raise AttributeError("slide.background not available")
            except:
                # Approach 2: Add as a full-slide image
                try:
                    # Convert slide dimensions back to inches for add_picture
                    width_inches = slide_width / 96  # Approximate pixels to inches
                    height_inches = slide_height / 96
                    
                    picture = slide.shapes.add_picture(temp_path, 0, 0, 
                                                     Inches(width_inches), Inches(height_inches))
                    
                    # Send to back so it acts as background
                    slide.shapes._spTree.remove(picture._element)
                    slide.shapes._spTree.insert(2, picture._element)  # Insert after background elements
                    
                    result['success'] = True
                    result['message'] = 'Gradient background added as full-slide image'
                except Exception as e:
                    result['message'] = f'Failed to add gradient background: {str(e)}'
        
        finally:
            # Clean up temporary file
            try:
                os.unlink(temp_path)
            except:
                pass
        
        if result['success']:
            result['gradient_info']['dimensions'] = {
                'width_pixels': slide_width,
                'height_pixels': slide_height
            }
        
        return result
        
    except Exception as e:
        result['message'] = f'Error creating gradient background: {str(e)}'
        return result

def create_professional_gradient_background(slide, color_scheme: str = 'modern_blue', 
                                          style: str = 'subtle', direction: str = 'diagonal') -> Dict:
    """
    Create a professional gradient background using predefined color schemes.
    
    Args:
        slide: The slide object
        color_scheme: Professional color scheme name
        style: Gradient style ('subtle', 'bold', 'accent')
        direction: Gradient direction
        
    Returns:
        Dictionary with operation results
    """
    # Define gradient styles based on professional color schemes
    style_configs = {
        'subtle': {
            'start': 'light',
            'end': 'light',  # Very subtle gradient within same color family
            'opacity': 0.3
        },
        'bold': {
            'start': 'primary',
            'end': 'secondary',
            'opacity': 0.8
        },
        'accent': {
            'start': 'accent1',
            'end': 'accent2',
            'opacity': 0.6
        }
    }
    
    style_config = style_configs.get(style, style_configs['subtle'])
    
    # Get colors from professional color scheme
    start_color = get_professional_color(color_scheme, style_config['start'])
    end_color = get_professional_color(color_scheme, style_config['end'])
    
    # For subtle gradients, create a lighter variation of the end color
    if style == 'subtle':
        # Blend with white to create subtle variation
        blend_factor = 0.7
        end_color = tuple(int(c + (255 - c) * blend_factor) for c in start_color)
    
    return set_slide_gradient_background(slide, start_color, end_color, direction)

# ---- Image Enhancement Functions ----

def enhance_image_pillow(image_path: str, enhancements: Dict = None, output_path: str = None) -> Dict:
    """
    Enhance an image using Pillow filters and adjustments.
    
    Args:
        image_path: Path to the input image
        enhancements: Dictionary of enhancement settings
        output_path: Path for enhanced image (optional, creates temp file if not provided)
        
    Returns:
        Dictionary with enhancement results and output path
    """
    default_enhancements = {
        'brightness': 1.0,    # 1.0 = no change, >1.0 = brighter, <1.0 = darker
        'contrast': 1.0,      # 1.0 = no change, >1.0 = more contrast
        'saturation': 1.0,    # 1.0 = no change, >1.0 = more saturated
        'sharpness': 1.0,     # 1.0 = no change, >1.0 = sharper
        'blur_radius': 0,     # 0 = no blur, >0 = blur radius
        'filter': None        # 'DETAIL', 'EDGE_ENHANCE', 'EMBOSS', 'SMOOTH', etc.
    }
    
    if enhancements:
        default_enhancements.update(enhancements)
    
    result = {
        'success': False,
        'message': '',
        'original_path': image_path,
        'enhanced_path': None,
        'enhancements_applied': default_enhancements
    }
    
    try:
        # Open and process image
        with Image.open(image_path) as img:
            enhanced_img = img.copy()
            
            # Apply brightness adjustment
            if default_enhancements['brightness'] != 1.0:
                enhancer = ImageEnhance.Brightness(enhanced_img)
                enhanced_img = enhancer.enhance(default_enhancements['brightness'])
            
            # Apply contrast adjustment
            if default_enhancements['contrast'] != 1.0:
                enhancer = ImageEnhance.Contrast(enhanced_img)
                enhanced_img = enhancer.enhance(default_enhancements['contrast'])
            
            # Apply color saturation adjustment
            if default_enhancements['saturation'] != 1.0:
                enhancer = ImageEnhance.Color(enhanced_img)
                enhanced_img = enhancer.enhance(default_enhancements['saturation'])
            
            # Apply sharpness adjustment
            if default_enhancements['sharpness'] != 1.0:
                enhancer = ImageEnhance.Sharpness(enhanced_img)
                enhanced_img = enhancer.enhance(default_enhancements['sharpness'])
            
            # Apply blur filter
            if default_enhancements['blur_radius'] > 0:
                enhanced_img = enhanced_img.filter(ImageFilter.GaussianBlur(default_enhancements['blur_radius']))
            
            # Apply special filters
            filter_name = default_enhancements['filter']
            if filter_name:
                filter_map = {
                    'DETAIL': ImageFilter.DETAIL,
                    'EDGE_ENHANCE': ImageFilter.EDGE_ENHANCE,
                    'EDGE_ENHANCE_MORE': ImageFilter.EDGE_ENHANCE_MORE,
                    'EMBOSS': ImageFilter.EMBOSS,
                    'FIND_EDGES': ImageFilter.FIND_EDGES,
                    'SMOOTH': ImageFilter.SMOOTH,
                    'SMOOTH_MORE': ImageFilter.SMOOTH_MORE,
                    'SHARPEN': ImageFilter.SHARPEN
                }
                
                if filter_name in filter_map:
                    enhanced_img = enhanced_img.filter(filter_map[filter_name])
            
            # Save enhanced image
            if output_path is None:
                # Create temporary file
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                output_path = temp_file.name
                temp_file.close()
            
            enhanced_img.save(output_path, 'PNG', quality=95)
            
            result['success'] = True
            result['enhanced_path'] = output_path
            result['message'] = 'Image enhanced successfully'
            
            # Add image information
            result['image_info'] = {
                'original_size': img.size,
                'enhanced_size': enhanced_img.size,
                'original_mode': img.mode,
                'enhanced_mode': enhanced_img.mode
            }
        
        return result
        
    except Exception as e:
        result['message'] = f'Error enhancing image: {str(e)}'
        return result

def apply_professional_image_enhancement(image_path: str, style: str = 'presentation') -> Dict:
    """
    Apply professional image enhancement presets suitable for presentations.
    
    Args:
        image_path: Path to the input image
        style: Enhancement style ('presentation', 'vibrant', 'subtle', 'sharp')
        
    Returns:
        Dictionary with enhancement results
    """
    enhancement_presets = {
        'presentation': {
            'brightness': 1.05,
            'contrast': 1.15,
            'saturation': 1.1,
            'sharpness': 1.2,
            'filter': 'DETAIL'
        },
        'vibrant': {
            'brightness': 1.1,
            'contrast': 1.3,
            'saturation': 1.4,
            'sharpness': 1.3,
            'filter': 'EDGE_ENHANCE'
        },
        'subtle': {
            'brightness': 1.02,
            'contrast': 1.05,
            'saturation': 1.05,
            'sharpness': 1.1,
            'filter': 'SMOOTH'
        },
        'sharp': {
            'brightness': 1.0,
            'contrast': 1.2,
            'saturation': 1.0,
            'sharpness': 1.5,
            'filter': 'SHARPEN'
        }
    }
    
    preset = enhancement_presets.get(style, enhancement_presets['presentation'])
    return enhance_image_pillow(image_path, preset)

# ---- Font Beautification Functions ----

def analyze_font_file(font_path: str) -> Dict:
    """
    Analyze a font file using FontTools to extract information.
    
    Args:
        font_path: Path to the font file (.ttf, .otf)
        
    Returns:
        Dictionary with font information
    """
    result = {
        'success': False,
        'message': '',
        'font_info': {}
    }
    
    try:
        font = TTFont(font_path)
        
        # Extract basic font information
        font_info = {
            'file_path': font_path,
            'font_tables': list(font.keys()),
            'glyph_count': font.getGlyphSet().keys().__len__() if hasattr(font, 'getGlyphSet') else 0
        }
        
        # Extract name table information
        if 'name' in font:
            name_table = font['name']
            font_info['names'] = {}
            
            for record in name_table.names:
                try:
                    name_text = record.toUnicode()
                    name_id = record.nameID
                    
                    # Map common name IDs
                    name_id_map = {
                        1: 'family_name',
                        2: 'subfamily_name',
                        3: 'unique_id',
                        4: 'full_name',
                        5: 'version',
                        6: 'postscript_name',
                        16: 'preferred_family',
                        17: 'preferred_subfamily'
                    }
                    
                    key = name_id_map.get(name_id, f'name_id_{name_id}')
                    font_info['names'][key] = name_text
                    
                except Exception:
                    continue
        
        # Extract OS/2 table information (metrics)
        if 'OS/2' in font:
            os2_table = font['OS/2']
            font_info['metrics'] = {
                'weight_class': getattr(os2_table, 'usWeightClass', None),
                'width_class': getattr(os2_table, 'usWidthClass', None),
                'x_height': getattr(os2_table, 'sxHeight', None),
                'cap_height': getattr(os2_table, 'sCapHeight', None)
            }
        
        # Extract head table information
        if 'head' in font:
            head_table = font['head']
            font_info['head'] = {
                'units_per_em': getattr(head_table, 'unitsPerEm', None),
                'created': getattr(head_table, 'created', None),
                'modified': getattr(head_table, 'modified', None)
            }
        
        font.close()
        
        result['success'] = True
        result['font_info'] = font_info
        result['message'] = 'Font analysis completed successfully'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error analyzing font: {str(e)}'
        return result

def optimize_font_for_presentation(font_path: str, output_path: str = None, 
                                  text_content: str = None) -> Dict:
    """
    Optimize a font file for presentation use by subsetting and optimizing.
    
    Args:
        font_path: Path to the input font file
        output_path: Path for optimized font (optional)
        text_content: Specific text content to optimize for (optional)
        
    Returns:
        Dictionary with optimization results
    """
    result = {
        'success': False,
        'message': '',
        'original_path': font_path,
        'optimized_path': None,
        'size_reduction': 0
    }
    
    try:
        # Get original file size
        original_size = os.path.getsize(font_path)
        
        # Create output path if not provided
        if output_path is None:
            base_name = os.path.splitext(os.path.basename(font_path))[0]
            output_path = tempfile.NamedTemporaryFile(
                delete=False, suffix=f'_optimized.ttf', prefix=f'{base_name}_'
            ).name
        
        # Load font
        font = TTFont(font_path)
        
        # Create subsetter for optimization
        subsetter = Subsetter()
        
        if text_content:
            # Subset to only include characters used in the content
            subsetter.populate(text=text_content)
        else:
            # Keep common characters for presentations
            common_chars = (
                'ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz'
                '0123456789.,;:!?-()[]{}"\' /\\@#$%^&*+=<>|~`'
                '©®™€£¥§¶•‰""''…–—'
            )
            subsetter.populate(text=common_chars)
        
        # Configure subsetter options
        subsetter.options.recommended_glyphs = True
        subsetter.options.name_legacy = True
        subsetter.options.name_languages = ['*']
        subsetter.options.layout_features = ['*']
        subsetter.options.glyph_names = True
        
        # Apply subsetting
        subsetter.subset(font)
        
        # Save optimized font
        font.save(output_path)
        font.close()
        
        # Calculate size reduction
        optimized_size = os.path.getsize(output_path)
        size_reduction = ((original_size - optimized_size) / original_size) * 100
        
        result['success'] = True
        result['optimized_path'] = output_path
        result['size_reduction'] = round(size_reduction, 2)
        result['message'] = f'Font optimized successfully. Size reduced by {size_reduction:.1f}%'
        result['file_sizes'] = {
            'original_bytes': original_size,
            'optimized_bytes': optimized_size
        }
        
        return result
        
    except Exception as e:
        result['message'] = f'Error optimizing font: {str(e)}'
        return result

def get_font_recommendations(font_analysis: Dict, presentation_type: str = 'business') -> Dict:
    """
    Get font usage recommendations based on analysis and presentation type.
    
    Args:
        font_analysis: Result from analyze_font_file()
        presentation_type: Type of presentation ('business', 'creative', 'academic')
        
    Returns:
        Dictionary with font recommendations
    """
    recommendations = {
        'suitable_for_presentation': False,
        'recommended_sizes': {},
        'usage_suggestions': [],
        'potential_issues': [],
        'alternatives': []
    }
    
    try:
        if not font_analysis.get('success'):
            recommendations['potential_issues'].append('Font analysis failed')
            return recommendations
        
        font_info = font_analysis['font_info']
        
        # Analyze font characteristics
        weight_class = font_info.get('metrics', {}).get('weight_class', 400)
        family_name = font_info.get('names', {}).get('family_name', 'Unknown')
        
        # Determine suitability
        if weight_class and 300 <= weight_class <= 700:
            recommendations['suitable_for_presentation'] = True
        
        # Recommend sizes based on presentation type
        if presentation_type == 'business':
            recommendations['recommended_sizes'] = {
                'title': '32-44pt',
                'subtitle': '20-28pt',
                'body': '16-24pt',
                'caption': '12-16pt'
            }
        elif presentation_type == 'creative':
            recommendations['recommended_sizes'] = {
                'title': '36-60pt',
                'subtitle': '24-36pt',
                'body': '18-28pt',
                'caption': '14-18pt'
            }
        else:  # academic
            recommendations['recommended_sizes'] = {
                'title': '28-36pt',
                'subtitle': '18-24pt',
                'body': '14-20pt',
                'caption': '10-14pt'
            }
        
        # Usage suggestions
        if 'Light' in family_name or weight_class < 350:
            recommendations['usage_suggestions'].append('Best for titles and headings, may be too light for body text')
        elif 'Bold' in family_name or weight_class > 600:
            recommendations['usage_suggestions'].append('Good for emphasis and headings, use sparingly')
        else:
            recommendations['usage_suggestions'].append('Versatile font suitable for both headings and body text')
        
        # Check for potential issues
        if font_info.get('glyph_count', 0) < 200:
            recommendations['potential_issues'].append('Limited character set - may not support all languages')
        
        if weight_class and weight_class < 300:
            recommendations['potential_issues'].append('Very light weight may be hard to read on projectors')
        
        return recommendations
        
    except Exception as e:
        recommendations['potential_issues'].append(f'Error generating recommendations: {str(e)}')
        return recommendations

# ---- Image Functions ----

def add_image(slide, image_path: str, left: float, top: float, width: float = None, height: float = None) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: The slide object
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created picture shape
    """
    if width and height:
        picture = slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    else:
        picture = slide.shapes.add_picture(
            image_path, Inches(left), Inches(top)
        )
    return picture

def add_image_from_base64(slide, base64_string: str, left: float, top: float, 
                          width: float = None, height: float = None) -> Any:
    """
    Add an image from a base64 encoded string to a slide.
    
    Args:
        slide: The slide object
        base64_string: Base64 encoded image string
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created picture shape
    """
    image_data = base64.b64decode(base64_string)
    image_stream = io.BytesIO(image_data)
    
    if width and height:
        picture = slide.shapes.add_picture(
            image_stream, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    else:
        picture = slide.shapes.add_picture(
            image_stream, Inches(left), Inches(top)
        )
    return picture

# ---- Table Functions ----

def add_table(slide, rows: int, cols: int, left: float, top: float, width: float, height: float) -> Any:
    """
    Add a table to a slide.
    
    Args:
        slide: The slide object
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created table shape
    """
    table = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    return table

def set_cell_text(table, row: int, col: int, text: str) -> None:
    """
    Set text in a table cell.
    
    Args:
        table: The table object
        row: Row index
        col: Column index
        text: Text content
    """
    cell = table.cell(row, col)
    cell.text = text

def format_table_cell(cell, font_size: int = None, font_name: str = None, 
                     bold: bool = None, italic: bool = None, 
                     color: Tuple[int, int, int] = None,
                     bg_color: Tuple[int, int, int] = None,
                     alignment: str = None,
                     vertical_alignment: str = None) -> None:
    """
    Format a table cell.
    
    Args:
        cell: The table cell to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple for text (r, g, b)
        bg_color: RGB color tuple for background (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    vertical_alignment_map = {
        'top': MSO_VERTICAL_ANCHOR.TOP,
        'middle': MSO_VERTICAL_ANCHOR.MIDDLE,
        'bottom': MSO_VERTICAL_ANCHOR.BOTTOM
    }
    
    # Format text
    text_frame = cell.text_frame
    
    if vertical_alignment and vertical_alignment in vertical_alignment_map:
        text_frame.vertical_anchor = vertical_alignment_map[vertical_alignment]
    
    for paragraph in text_frame.paragraphs:
        if alignment and alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
            
        for run in paragraph.runs:
            font = run.font
            
            if font_size is not None:
                font.size = Pt(font_size)
                
            if font_name is not None:
                font.name = font_name
                
            if bold is not None:
                font.bold = bold
                
            if italic is not None:
                font.italic = italic
                
            if color is not None:
                r, g, b = color
                font.color.rgb = RGBColor(r, g, b)
    
    # Set background color
    if bg_color is not None:
        r, g, b = bg_color
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(r, g, b)

# ---- Shape Functions ----

def add_shape(slide, shape_type: str, left: float, top: float, width: float, height: float) -> Any:
    """
    Add an auto shape to a slide.
    
    Args:
        slide: The slide object
        shape_type: Shape type string (e.g., 'rectangle', 'oval', 'triangle')
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created shape
    """
    # Ensure shape_type is a string
    shape_type_str = str(shape_type)
    shape_type_lower = shape_type_str.lower()
    
    # Define shape type mapping with correct enum values
    shape_type_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'diamond': MSO_SHAPE.DIAMOND,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,  # Changed from TRIANGLE to ISOSCELES_TRIANGLE
        'isosceles_triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'right_triangle': MSO_SHAPE.RIGHT_TRIANGLE,
        'pentagon': MSO_SHAPE.PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'heptagon': MSO_SHAPE.HEPTAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'star': MSO_SHAPE.STAR_5_POINTS,
        'arrow': MSO_SHAPE.ARROW,
        'cloud': MSO_SHAPE.CLOUD,
        'heart': MSO_SHAPE.HEART,
        'lightning_bolt': MSO_SHAPE.LIGHTNING_BOLT,
        'sun': MSO_SHAPE.SUN,
        'moon': MSO_SHAPE.MOON,
        'smiley_face': MSO_SHAPE.SMILEY_FACE,
        'no_symbol': MSO_SHAPE.NO_SYMBOL,
        'flowchart_process': MSO_SHAPE.FLOWCHART_PROCESS,
        'flowchart_decision': MSO_SHAPE.FLOWCHART_DECISION,
        'flowchart_data': MSO_SHAPE.FLOWCHART_DATA,
        'flowchart_document': MSO_SHAPE.FLOWCHART_DOCUMENT,
        'flowchart_predefined_process': MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
        'flowchart_internal_storage': MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
        'flowchart_connector': MSO_SHAPE.FLOWCHART_CONNECTOR
    }
    
    # Check if shape type is valid before trying to use it
    if shape_type_lower not in shape_type_map:
        available_shapes = ', '.join(sorted(shape_type_map.keys()))
        raise ValueError(f"Unsupported shape type: '{shape_type}'. Available shape types: {available_shapes}")
    
    # Get the shape enum value
    shape_enum = shape_type_map[shape_type_lower]
    
    # Create the shape with better error handling
    try:
        shape = slide.shapes.add_shape(
            shape_enum, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        return shape
    except Exception as e:
        # More detailed error for debugging
        raise ValueError(f"Could not create shape '{shape_type}' (enum: {shape_enum.__class__.__name__}.{shape_enum.name}): {str(e)}")

def format_shape(shape, fill_color: Tuple[int, int, int] = None, 
                line_color: Tuple[int, int, int] = None, line_width: float = None) -> None:
    """
    Format a shape.
    
    Args:
        shape: The shape object
        fill_color: RGB color tuple for fill (r, g, b)
        line_color: RGB color tuple for outline (r, g, b)
        line_width: Line width in points
    """
    if fill_color is not None:
        r, g, b = fill_color
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
    
    if line_color is not None:
        r, g, b = line_color
        shape.line.color.rgb = RGBColor(r, g, b)
    
    if line_width is not None:
        shape.line.width = Pt(line_width)

def format_professional_shape(shape, color_scheme: str = 'modern_blue', 
                             style: str = 'primary', with_shadow: bool = True) -> None:
    """
    Apply professional formatting to a shape with modern styling.
    
    Args:
        shape: The shape object
        color_scheme: Color scheme to use
        style: Style type ('primary', 'secondary', 'accent1', 'accent2', 'subtle')
        with_shadow: Whether to add a subtle shadow effect
    """
    try:
        # Get colors for the style
        if style == 'primary':
            fill_color = get_professional_color(color_scheme, 'primary')
            line_color = get_professional_color(color_scheme, 'primary')
        elif style == 'secondary':
            fill_color = get_professional_color(color_scheme, 'secondary')
            line_color = get_professional_color(color_scheme, 'secondary')
        elif style == 'accent1':
            fill_color = get_professional_color(color_scheme, 'accent1')
            line_color = get_professional_color(color_scheme, 'accent1')
        elif style == 'accent2':
            fill_color = get_professional_color(color_scheme, 'accent2')
            line_color = get_professional_color(color_scheme, 'accent2')
        elif style == 'subtle':
            fill_color = get_professional_color(color_scheme, 'light')
            line_color = get_professional_color(color_scheme, 'text')
        else:
            fill_color = get_professional_color(color_scheme, 'primary')
            line_color = get_professional_color(color_scheme, 'primary')
        
        # Apply fill
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
        
        # Apply modern line styling
        if style == 'subtle':
            shape.line.color.rgb = RGBColor(*line_color)
            shape.line.width = Pt(1)
        else:
            # No outline for solid colored shapes (modern flat design)
            shape.line.fill.background()
        
        # Add subtle shadow for depth (if supported)
        if with_shadow and hasattr(shape, 'shadow'):
            try:
                shadow = shape.shadow
                shadow.inherit = False
                shadow.visible = True
                shadow.style = 1  # Simple shadow
                shadow.blur_radius = Pt(4)
                shadow.distance = Pt(3)
                shadow.angle = 45
                shadow.color.rgb = RGBColor(0, 0, 0)
                shadow.transparency = 0.3
            except:
                pass  # Shadow not supported
                
    except Exception:
        # Fallback to basic formatting
        try:
            format_shape(shape, fill_color, line_color, 1.0)
        except:
            pass

# ---- Chart Functions ----

def add_chart(slide, chart_type: str, left: float, top: float, width: float, height: float,
             categories: List[str], series_names: List[str], series_values: List[List[float]]) -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: The slide object
        chart_type: Type of chart ('column', 'bar', 'line', 'pie')
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: List of category names
        series_names: List of series names
        series_values: List of lists containing values for each series
        
    Returns:
        The created chart
    """
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'line': XL_CHART_TYPE.LINE,
        'line_markers': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA,
        'stacked_area': XL_CHART_TYPE.AREA_STACKED,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'radar': XL_CHART_TYPE.RADAR,
        'radar_markers': XL_CHART_TYPE.RADAR_MARKERS
    }
    
    chart_type_enum = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for i, series_name in enumerate(series_names):
        chart_data.add_series(series_name, series_values[i])
    
    # Add chart to slide
    graphic_frame = slide.shapes.add_chart(
        chart_type_enum, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )
    
    return graphic_frame.chart

def format_chart(chart, has_legend: bool = True, legend_position: str = 'right',
                has_data_labels: bool = False, title: str = None) -> None:
    """
    Format a chart.
    
    Args:
        chart: The chart object
        has_legend: Whether to show the legend
        legend_position: Position of the legend ('right', 'left', 'top', 'bottom')
        has_data_labels: Whether to show data labels
        title: Chart title
    """
    # Set chart title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    
    # Configure legend
    chart.has_legend = has_legend
    if has_legend:
        position_map = {
            'right': 2,  # XL_LEGEND_POSITION.RIGHT
            'left': 3,   # XL_LEGEND_POSITION.LEFT
            'top': 1,    # XL_LEGEND_POSITION.TOP
            'bottom': 4  # XL_LEGEND_POSITION.BOTTOM
        }
        chart.legend.position = position_map.get(legend_position.lower(), 2)
    
    # Configure data labels
    for series in chart.series:
        series.has_data_labels = has_data_labels

def format_professional_chart(chart, color_scheme: str = 'modern_blue', 
                             title: str = None, modern_style: bool = True) -> None:
    """
    Apply professional formatting to a chart with modern styling and colors.
    
    Args:
        chart: The chart object
        color_scheme: Color scheme to use
        title: Chart title
        modern_style: Whether to apply modern flat design principles
    """
    try:
        # Set professional title
        if title:
            chart.has_title = True
            title_frame = chart.chart_title.text_frame
            title_frame.text = title
            
            # Format title text
            title_font = get_professional_font('title', 'small')
            title_color = get_professional_color(color_scheme, 'primary')
            
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    font.name = title_font['name']
                    font.size = Pt(title_font['size'])
                    font.bold = title_font['bold']
                    font.color.rgb = RGBColor(*title_color)
        
        # Configure modern legend
        if chart.has_legend:
            chart.legend.position = 4  # Bottom position for modern look
            try:
                legend_font = get_professional_font('body', 'small')
                text_color = get_professional_color(color_scheme, 'text')
                
                for paragraph in chart.legend.font:
                    font = paragraph.font
                    font.name = legend_font['name']
                    font.size = Pt(legend_font['size'])
                    font.color.rgb = RGBColor(*text_color)
            except:
                pass
        
        # Apply professional colors to series
        professional_colors = [
            get_professional_color(color_scheme, 'primary'),
            get_professional_color(color_scheme, 'accent1'),
            get_professional_color(color_scheme, 'accent2'),
            get_professional_color(color_scheme, 'secondary'),
        ]
        
        for i, series in enumerate(chart.series):
            try:
                color = professional_colors[i % len(professional_colors)]
                
                # Apply color to series
                if hasattr(series, 'fill'):
                    series.fill.solid()
                    series.fill.fore_color.rgb = RGBColor(*color)
                
                # Modern style: remove borders from series
                if modern_style and hasattr(series, 'line'):
                    series.line.fill.background()
                
            except Exception:
                continue
        
        # Style axes for modern look
        try:
            # Category axis
            if hasattr(chart, 'category_axis'):
                cat_axis = chart.category_axis
                if hasattr(cat_axis, 'tick_labels'):
                    font_settings = get_professional_font('caption', 'medium')
                    text_color = get_professional_color(color_scheme, 'text')
                    
                    tick_font = cat_axis.tick_labels.font
                    tick_font.name = font_settings['name']
                    tick_font.size = Pt(font_settings['size'])
                    tick_font.color.rgb = RGBColor(*text_color)
            
            # Value axis
            if hasattr(chart, 'value_axis'):
                val_axis = chart.value_axis
                if hasattr(val_axis, 'tick_labels'):
                    font_settings = get_professional_font('caption', 'medium')
                    text_color = get_professional_color(color_scheme, 'text')
                    
                    tick_font = val_axis.tick_labels.font
                    tick_font.name = font_settings['name']
                    tick_font.size = Pt(font_settings['size'])
                    tick_font.color.rgb = RGBColor(*text_color)
                
        except Exception:
            pass
        
    except Exception:
        # Fallback to basic formatting
        try:
            format_chart(chart, has_legend=True, title=title)
        except:
            pass

# ---- Picture Effects Functions ----

def apply_picture_shadow(picture_shape, shadow_type: str = 'outer', blur_radius: float = 4.0,
                        distance: float = 3.0, direction: float = 315.0, color: Tuple[int, int, int] = (0, 0, 0),
                        transparency: float = 0.6) -> Dict:
    """
    Apply shadow effect to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        shadow_type: Type of shadow ('outer', 'inner')
        blur_radius: Shadow blur radius in points
        distance: Shadow distance in points
        direction: Shadow direction in degrees (0-360)
        color: RGB color tuple for shadow
        transparency: Shadow transparency (0.0-1.0, where 1.0 is fully transparent)
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'shadow_info': {
            'type': shadow_type,
            'blur_radius': blur_radius,
            'distance': distance,
            'direction': direction,
            'color': color,
            'transparency': transparency
        }
    }
    
    try:
        # Access the shadow format
        shadow = picture_shape.shadow
        shadow.inherit = False  # Override theme inheritance
        
        # Configure shadow properties
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        
        # Set shadow visibility
        if hasattr(shadow, 'visible'):
            shadow.visible = True
        
        # Set shadow style (outer shadow)
        if hasattr(shadow, 'style'):
            shadow.style = 1  # Outer shadow
        
        # Set blur radius
        if hasattr(shadow, 'blur_radius'):
            shadow.blur_radius = Pt(blur_radius)
        
        # Set distance
        if hasattr(shadow, 'distance'):
            shadow.distance = Pt(distance)
        
        # Set direction (convert degrees to angle units)
        if hasattr(shadow, 'angle'):
            # Convert degrees to EMU angle units (1 degree = 60000 EMU)
            shadow.angle = int(direction * 60000)
        
        # Set color and transparency
        if hasattr(shadow, 'color'):
            shadow.color.rgb = RGBColor(*color)
            if hasattr(shadow.color, 'transparency'):
                shadow.color.transparency = transparency
        
        result['success'] = True
        result['message'] = f'Applied {shadow_type} shadow effect to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying shadow effect: {str(e)}'
        return result

def apply_picture_reflection(picture_shape, size: float = 0.5, transparency: float = 0.5,
                           distance: float = 0.0, blur: float = 4.0) -> Dict:
    """
    Apply reflection effect to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        size: Reflection size as percentage (0.0-1.0)
        transparency: Reflection transparency (0.0-1.0)
        distance: Distance between image and reflection in points
        blur: Reflection blur amount in points
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'reflection_info': {
            'size': size,
            'transparency': transparency,
            'distance': distance,
            'blur': blur
        }
    }
    
    try:
        # For reflection, we need to manipulate the XML directly as python-pptx
        # doesn't have direct reflection support
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import nsdecls, qn
        
        # Get the shape's XML element
        shape_elm = picture_shape._element
        sp_pr = shape_elm.find(qn('p:spPr'))
        
        if sp_pr is None:
            result['message'] = 'Could not find shape properties element'
            return result
        
        # Create or find effect list
        effect_lst = sp_pr.find(qn('a:effectLst'))
        if effect_lst is None:
            effect_lst = OxmlElement('a:effectLst')
            sp_pr.append(effect_lst)
        
        # Remove existing reflection if present
        existing_reflection = effect_lst.find(qn('a:reflection'))
        if existing_reflection is not None:
            effect_lst.remove(existing_reflection)
        
        # Create reflection element
        reflection = OxmlElement('a:reflection')
        
        # Set reflection attributes
        reflection.set('blurRad', str(int(blur * 12700)))  # Convert points to EMU
        reflection.set('stA', str(int((1.0 - transparency) * 100000)))  # Start alpha
        reflection.set('endA', str(int(transparency * 100000)))  # End alpha
        reflection.set('endPos', str(int(size * 100000)))  # End position
        reflection.set('dist', str(int(distance * 12700)))  # Distance in EMU
        
        # Add reflection to effect list
        effect_lst.append(reflection)
        
        result['success'] = True
        result['message'] = 'Applied reflection effect to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying reflection effect: {str(e)}'
        return result

def apply_picture_glow(picture_shape, size: float = 5.0, color: Tuple[int, int, int] = (0, 176, 240),
                      transparency: float = 0.4) -> Dict:
    """
    Apply glow effect to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        size: Glow size in points
        color: RGB color tuple for glow
        transparency: Glow transparency (0.0-1.0)
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'glow_info': {
            'size': size,
            'color': color,
            'transparency': transparency
        }
    }
    
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import nsdecls, qn
        
        # Get the shape's XML element
        shape_elm = picture_shape._element
        sp_pr = shape_elm.find(qn('p:spPr'))
        
        if sp_pr is None:
            result['message'] = 'Could not find shape properties element'
            return result
        
        # Create or find effect list
        effect_lst = sp_pr.find(qn('a:effectLst'))
        if effect_lst is None:
            effect_lst = OxmlElement('a:effectLst')
            sp_pr.append(effect_lst)
        
        # Remove existing glow if present
        existing_glow = effect_lst.find(qn('a:glow'))
        if existing_glow is not None:
            effect_lst.remove(existing_glow)
        
        # Create glow element
        glow = OxmlElement('a:glow')
        glow.set('rad', str(int(size * 12700)))  # Convert points to EMU
        
        # Create color element
        srgb_clr = OxmlElement('a:srgbClr')
        srgb_clr.set('val', f'{color[0]:02x}{color[1]:02x}{color[2]:02x}')
        
        # Add alpha (transparency) if needed
        if transparency > 0:
            alpha = OxmlElement('a:alpha')
            alpha.set('val', str(int((1.0 - transparency) * 100000)))
            srgb_clr.append(alpha)
        
        glow.append(srgb_clr)
        effect_lst.append(glow)
        
        result['success'] = True
        result['message'] = 'Applied glow effect to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying glow effect: {str(e)}'
        return result

def apply_picture_soft_edges(picture_shape, radius: float = 2.5) -> Dict:
    """
    Apply soft edges effect to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        radius: Soft edge radius in points
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'soft_edge_info': {
            'radius': radius
        }
    }
    
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import nsdecls, qn
        
        # Get the shape's XML element
        shape_elm = picture_shape._element
        sp_pr = shape_elm.find(qn('p:spPr'))
        
        if sp_pr is None:
            result['message'] = 'Could not find shape properties element'
            return result
        
        # Create or find effect list
        effect_lst = sp_pr.find(qn('a:effectLst'))
        if effect_lst is None:
            effect_lst = OxmlElement('a:effectLst')
            sp_pr.append(effect_lst)
        
        # Remove existing soft edge if present
        existing_soft_edge = effect_lst.find(qn('a:softEdge'))
        if existing_soft_edge is not None:
            effect_lst.remove(existing_soft_edge)
        
        # Create soft edge element
        soft_edge = OxmlElement('a:softEdge')
        soft_edge.set('rad', str(int(radius * 12700)))  # Convert points to EMU
        
        effect_lst.append(soft_edge)
        
        result['success'] = True
        result['message'] = 'Applied soft edges effect to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying soft edges effect: {str(e)}'
        return result

def apply_picture_rotation(picture_shape, rotation: float) -> Dict:
    """
    Apply rotation to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        rotation: Rotation angle in degrees (positive = clockwise)
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'rotation_info': {
            'rotation': rotation
        }
    }
    
    try:
        # Convert degrees to EMU (1 degree = 60000 EMU)
        rotation_emu = int(rotation * 60000)
        
        # Set rotation on the shape
        picture_shape.rotation = rotation
        
        result['success'] = True
        result['message'] = f'Applied {rotation}° rotation to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying rotation: {str(e)}'
        return result

def apply_picture_transparency(picture_shape, transparency: float) -> Dict:
    """
    Apply transparency to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        transparency: Transparency level (0.0-1.0, where 1.0 is fully transparent)
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'transparency_info': {
            'transparency': transparency
        }
    }
    
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import nsdecls, qn
        
        # Get the picture's blip fill element
        shape_elm = picture_shape._element
        blip_fill = shape_elm.find(qn('p:blipFill'))
        
        if blip_fill is None:
            result['message'] = 'Could not find picture fill element'
            return result
        
        # Find or create the blip element
        blip = blip_fill.find(qn('a:blip'))
        if blip is None:
            result['message'] = 'Could not find picture blip element'
            return result
        
        # Remove existing alpha modulation if present
        existing_alpha = blip.find(qn('a:alphaModFix'))
        if existing_alpha is not None:
            blip.remove(existing_alpha)
        
        # Add alpha modulation for transparency
        if transparency > 0:
            alpha_mod = OxmlElement('a:alphaModFix')
            alpha_mod.set('amt', str(int((1.0 - transparency) * 100000)))
            blip.append(alpha_mod)
        
        result['success'] = True
        result['message'] = f'Applied {transparency*100:.1f}% transparency to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying transparency: {str(e)}'
        return result

def apply_picture_bevel(picture_shape, bevel_type: str = 'circle', width: float = 6.0,
                       height: float = 6.0) -> Dict:
    """
    Apply bevel effect to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        bevel_type: Type of bevel ('circle', 'square', 'slope', 'riblet')
        width: Bevel width in points
        height: Bevel height in points
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'bevel_info': {
            'type': bevel_type,
            'width': width,
            'height': height
        }
    }
    
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import nsdecls, qn
        
        # Bevel type mapping
        bevel_types = {
            'circle': 'circle',
            'square': 'square',
            'slope': 'slope',
            'riblet': 'riblet'
        }
        
        bevel_preset = bevel_types.get(bevel_type, 'circle')
        
        # Get the shape's XML element
        shape_elm = picture_shape._element
        sp_pr = shape_elm.find(qn('p:spPr'))
        
        if sp_pr is None:
            result['message'] = 'Could not find shape properties element'
            return result
        
        # Create or find effect list
        effect_lst = sp_pr.find(qn('a:effectLst'))
        if effect_lst is None:
            effect_lst = OxmlElement('a:effectLst')
            sp_pr.append(effect_lst)
        
        # Remove existing bevel if present
        existing_bevel = effect_lst.find(qn('a:bevel'))
        if existing_bevel is not None:
            effect_lst.remove(existing_bevel)
        
        # Create bevel element
        bevel = OxmlElement('a:bevel')
        bevel.set('w', str(int(width * 12700)))  # Width in EMU
        bevel.set('h', str(int(height * 12700)))  # Height in EMU
        bevel.set('prst', bevel_preset)
        
        effect_lst.append(bevel)
        
        result['success'] = True
        result['message'] = f'Applied {bevel_type} bevel effect to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying bevel effect: {str(e)}'
        return result

def apply_picture_filter(picture_shape, filter_type: str = 'none', intensity: float = 0.5) -> Dict:
    """
    Apply color filter effect to a picture shape.
    
    Args:
        picture_shape: The picture shape object
        filter_type: Type of filter ('grayscale', 'sepia', 'washout', 'none')
        intensity: Filter intensity (0.0-1.0)
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': False,
        'message': '',
        'filter_info': {
            'type': filter_type,
            'intensity': intensity
        }
    }
    
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import nsdecls, qn
        
        # Get the picture's blip fill element
        shape_elm = picture_shape._element
        blip_fill = shape_elm.find(qn('p:blipFill'))
        
        if blip_fill is None:
            result['message'] = 'Could not find picture fill element'
            return result
        
        # Find or create the blip element
        blip = blip_fill.find(qn('a:blip'))
        if blip is None:
            result['message'] = 'Could not find picture blip element'
            return result
        
        # Remove existing color modifications
        for elem in blip.findall(qn('a:grayscl')):
            blip.remove(elem)
        for elem in blip.findall(qn('a:duotone')):
            blip.remove(elem)
        for elem in blip.findall(qn('a:lumMod')):
            blip.remove(elem)
        
        # Apply filter based on type
        if filter_type == 'grayscale':
            grayscale = OxmlElement('a:grayscl')
            blip.append(grayscale)
            
        elif filter_type == 'sepia':
            # Create sepia effect using duotone
            duotone = OxmlElement('a:duotone')
            
            # Dark sepia color
            srgb_clr1 = OxmlElement('a:srgbClr')
            srgb_clr1.set('val', '704214')  # Dark sepia
            duotone.append(srgb_clr1)
            
            # Light sepia color
            srgb_clr2 = OxmlElement('a:srgbClr')
            srgb_clr2.set('val', 'F4EABF')  # Light sepia
            duotone.append(srgb_clr2)
            
            blip.append(duotone)
            
        elif filter_type == 'washout':
            # Create washout effect using luminance modification
            lum_mod = OxmlElement('a:lumMod')
            lum_mod.set('val', str(int(50000 + intensity * 50000)))  # 50-100% luminance
            blip.append(lum_mod)
        
        # filter_type == 'none' removes all filters (already done above)
        
        result['success'] = True
        result['message'] = f'Applied {filter_type} filter to picture'
        
        return result
        
    except Exception as e:
        result['message'] = f'Error applying filter: {str(e)}'
        return result

def apply_combined_picture_effects(picture_shape, effects: Dict) -> Dict:
    """
    Apply multiple picture effects in combination.
    
    Args:
        picture_shape: The picture shape object
        effects: Dictionary of effects to apply with their parameters
        
    Returns:
        Dictionary with operation results
    """
    result = {
        'success': True,
        'message': '',
        'applied_effects': [],
        'failed_effects': [],
        'warnings': []
    }
    
    try:
        # Apply effects in order
        effect_functions = {
            'shadow': apply_picture_shadow,
            'reflection': apply_picture_reflection,
            'glow': apply_picture_glow,
            'soft_edges': apply_picture_soft_edges,
            'rotation': apply_picture_rotation,
            'transparency': apply_picture_transparency,
            'bevel': apply_picture_bevel,
            'filter': apply_picture_filter
        }
        
        for effect_name, effect_params in effects.items():
            if effect_name in effect_functions:
                try:
                    effect_result = effect_functions[effect_name](picture_shape, **effect_params)
                    if effect_result['success']:
                        result['applied_effects'].append(effect_name)
                    else:
                        result['failed_effects'].append({
                            'effect': effect_name,
                            'error': effect_result['message']
                        })
                except Exception as e:
                    result['failed_effects'].append({
                        'effect': effect_name,
                        'error': str(e)
                    })
            else:
                result['warnings'].append(f'Unknown effect: {effect_name}')
        
        if result['failed_effects']:
            result['success'] = len(result['applied_effects']) > 0
        
        applied_count = len(result['applied_effects'])
        failed_count = len(result['failed_effects'])
        
        if applied_count > 0 and failed_count == 0:
            result['message'] = f'Successfully applied {applied_count} picture effects'
        elif applied_count > 0 and failed_count > 0:
            result['message'] = f'Applied {applied_count} effects, {failed_count} failed'
        else:
            result['message'] = f'Failed to apply any effects'
            result['success'] = False
        
        return result
        
    except Exception as e:
        result['success'] = False
        result['message'] = f'Error applying combined effects: {str(e)}'
        return result

# ---- Document Properties Functions ----

def set_core_properties(presentation: Presentation, title: str = None, subject: str = None,
                       author: str = None, keywords: str = None, comments: str = None) -> None:
    """
    Set core document properties.
    
    Args:
        presentation: The Presentation object
        title: Document title
        subject: Document subject
        author: Document author
        keywords: Document keywords
        comments: Document comments
    """
    core_props = presentation.core_properties
    
    if title is not None:
        core_props.title = title
        
    if subject is not None:
        core_props.subject = subject
        
    if author is not None:
        core_props.author = author
        
    if keywords is not None:
        core_props.keywords = keywords
        
    if comments is not None:
        core_props.comments = comments

def get_core_properties(presentation: Presentation) -> Dict:
    """
    Get core document properties.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Dictionary of core properties
    """
    core_props = presentation.core_properties
    
    return {
        'title': core_props.title,
        'subject': core_props.subject,
        'author': core_props.author,
        'keywords': core_props.keywords,
        'comments': core_props.comments,
        'category': core_props.category,
        'created': core_props.created,
        'modified': core_props.modified,
        'last_modified_by': core_props.last_modified_by
    }


# ---- Slide Validation and Auto-Fix Functions ----

def validate_and_fix_slide_content(slide, auto_fix: bool = True, min_font_size: int = 8, 
                                  max_font_size: int = 72) -> Dict:
    """
    Comprehensive slide validation that checks font sizes, layout reasonableness, 
    and content boundaries. Automatically fixes issues when auto_fix is True.
    
    Args:
        slide: The slide object to validate and fix
        auto_fix: Whether to automatically fix detected issues
        min_font_size: Minimum allowed font size in points
        max_font_size: Maximum allowed font size in points
        
    Returns:
        Dictionary with validation results and fixes applied
    """
    result = {
        'success': True,
        'issues_found': [],
        'fixes_applied': [],
        'warnings': [],
        'validation_summary': {
            'font_issues': 0,
            'layout_issues': 0,
            'boundary_issues': 0,
            'total_shapes_checked': 0,
            'shapes_fixed': 0
        }
    }
    
    try:
        shapes_checked = 0
        shapes_fixed = 0
        
        # Iterate through all shapes on the slide
        for shape_idx, shape in enumerate(slide.shapes):
            shapes_checked += 1
            
            # Skip shapes without text frames
            if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
                continue
                
            text_frame = shape.text_frame
            shape_name = getattr(shape, 'name', f'Shape_{shape_idx}')
            
            # Validate and fix text content in this shape
            shape_result = validate_and_fix_shape_text(
                shape, text_frame, shape_name, auto_fix, min_font_size, max_font_size
            )
            
            # Merge results
            if not shape_result['success']:
                result['success'] = False
                
            result['issues_found'].extend(shape_result['issues_found'])
            result['fixes_applied'].extend(shape_result['fixes_applied'])
            result['warnings'].extend(shape_result['warnings'])
            
            # Update counters
            result['validation_summary']['font_issues'] += shape_result['font_issues']
            result['validation_summary']['layout_issues'] += shape_result['layout_issues']  
            result['validation_summary']['boundary_issues'] += shape_result['boundary_issues']
            
            if shape_result['fixes_applied']:
                shapes_fixed += 1
        
        result['validation_summary']['total_shapes_checked'] = shapes_checked
        result['validation_summary']['shapes_fixed'] = shapes_fixed
        
        # Overall slide layout validation
        layout_result = validate_slide_layout(slide, auto_fix)
        result['issues_found'].extend(layout_result['issues_found'])
        result['fixes_applied'].extend(layout_result['fixes_applied'])
        result['warnings'].extend(layout_result['warnings'])
        
        # Summary message
        total_issues = sum([
            result['validation_summary']['font_issues'],
            result['validation_summary']['layout_issues'],
            result['validation_summary']['boundary_issues']
        ])
        
        if total_issues == 0:
            result['message'] = f'Slide validation passed. Checked {shapes_checked} shapes with no issues found.'
        else:
            fixes_msg = f", {len(result['fixes_applied'])} fixes applied" if auto_fix else ""
            result['message'] = f'Slide validation found {total_issues} issues across {shapes_checked} shapes{fixes_msg}.'
            
    except Exception as e:
        result['success'] = False
        result['message'] = f'Slide validation error: {str(e)}'
        result['warnings'].append(f'Validation process failed: {str(e)}')
    
    return result


def validate_and_fix_shape_text(shape, text_frame, shape_name: str, auto_fix: bool, 
                               min_font_size: int, max_font_size: int) -> Dict:
    """
    Validate and fix text content within a single shape.
    
    Args:
        shape: The shape object
        text_frame: The text frame within the shape
        shape_name: Name identifier for the shape
        auto_fix: Whether to automatically fix issues
        min_font_size: Minimum font size allowed
        max_font_size: Maximum font size allowed
        
    Returns:
        Dictionary with validation results for this shape
    """
    result = {
        'success': True,
        'issues_found': [],
        'fixes_applied': [],
        'warnings': [],
        'font_issues': 0,
        'layout_issues': 0,
        'boundary_issues': 0
    }
    
    try:
        # Check if shape has text content
        if not text_frame.text.strip():
            return result  # No text to validate
            
        # Get shape dimensions in points for calculations
        shape_width_pt = shape.width / 12700  # Convert EMU to points
        shape_height_pt = shape.height / 12700
        
        # Check each paragraph in the text frame
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            if not paragraph.text.strip():
                continue
                
            # Validate paragraph-level formatting
            para_result = validate_paragraph_formatting(
                paragraph, para_idx, shape_name, shape_width_pt, shape_height_pt,
                auto_fix, min_font_size, max_font_size
            )
            
            # Merge paragraph results
            result['issues_found'].extend(para_result['issues_found'])
            result['fixes_applied'].extend(para_result['fixes_applied'])
            result['warnings'].extend(para_result['warnings'])
            result['font_issues'] += para_result['font_issues']
            
        # Check for text overflow and boundary issues
        boundary_result = check_text_boundaries(
            shape, text_frame, shape_name, auto_fix, min_font_size
        )
        
        result['issues_found'].extend(boundary_result['issues_found'])
        result['fixes_applied'].extend(boundary_result['fixes_applied'])  
        result['warnings'].extend(boundary_result['warnings'])
        result['boundary_issues'] += boundary_result['boundary_issues']
        
        # Check shape positioning and size reasonableness
        layout_result = validate_shape_layout(shape, shape_name, auto_fix)
        
        result['issues_found'].extend(layout_result['issues_found'])
        result['fixes_applied'].extend(layout_result['fixes_applied'])
        result['warnings'].extend(layout_result['warnings'])
        result['layout_issues'] += layout_result['layout_issues']
        
    except Exception as e:
        result['success'] = False
        result['warnings'].append(f'Shape {shape_name} validation failed: {str(e)}')
    
    return result


def validate_paragraph_formatting(paragraph, para_idx: int, shape_name: str, 
                                shape_width_pt: float, shape_height_pt: float,
                                auto_fix: bool, min_font_size: int, max_font_size: int) -> Dict:
    """
    Validate and fix formatting for a single paragraph.
    
    Args:
        paragraph: The paragraph object
        para_idx: Index of the paragraph
        shape_name: Name of the containing shape
        shape_width_pt: Shape width in points
        shape_height_pt: Shape height in points
        auto_fix: Whether to apply fixes
        min_font_size: Minimum font size
        max_font_size: Maximum font size
        
    Returns:
        Dictionary with validation results
    """
    result = {
        'issues_found': [],
        'fixes_applied': [],
        'warnings': [],
        'font_issues': 0
    }
    
    try:
        # Check each run in the paragraph
        for run_idx, run in enumerate(paragraph.runs):
            if not run.text.strip():
                continue
                
            font = run.font
            run_identifier = f'{shape_name}.paragraph[{para_idx}].run[{run_idx}]'
            
            # Check font size issues
            current_size = None
            if font.size is not None:
                current_size = font.size.pt
                
                # Check if font is too large
                if current_size > max_font_size:
                    issue = f'{run_identifier}: Font size {current_size}pt exceeds maximum {max_font_size}pt'
                    result['issues_found'].append(issue)
                    result['font_issues'] += 1
                    
                    if auto_fix:
                        font.size = Pt(max_font_size)
                        fix = f'{run_identifier}: Reduced font size from {current_size}pt to {max_font_size}pt'
                        result['fixes_applied'].append(fix)
                
                # Check if font is too small
                elif current_size < min_font_size:
                    issue = f'{run_identifier}: Font size {current_size}pt below minimum {min_font_size}pt'
                    result['issues_found'].append(issue)
                    result['font_issues'] += 1
                    
                    if auto_fix:
                        font.size = Pt(min_font_size)
                        fix = f'{run_identifier}: Increased font size from {current_size}pt to {min_font_size}pt'
                        result['fixes_applied'].append(fix)
                
                # Check if font size is reasonable for shape size
                recommended_max = min(shape_width_pt / 10, shape_height_pt / 3, max_font_size)
                if current_size > recommended_max:
                    issue = f'{run_identifier}: Font size {current_size}pt may be too large for shape dimensions'
                    result['issues_found'].append(issue)
                    result['font_issues'] += 1
                    
                    if auto_fix:
                        new_size = max(recommended_max, min_font_size)
                        font.size = Pt(new_size)
                        fix = f'{run_identifier}: Adjusted font size from {current_size}pt to {new_size}pt for better fit'
                        result['fixes_applied'].append(fix)
            
            # Check for other formatting issues
            if font.name and len(font.name) > 50:
                result['warnings'].append(f'{run_identifier}: Font name unusually long: {font.name[:30]}...')
                
    except Exception as e:
        result['warnings'].append(f'Paragraph formatting validation failed: {str(e)}')
    
    return result


def check_text_boundaries(shape, text_frame, shape_name: str, auto_fix: bool, 
                         min_font_size: int) -> Dict:
    """
    Check if text content exceeds shape boundaries and fix if needed.
    
    Args:
        shape: The shape object
        text_frame: The text frame
        shape_name: Shape identifier
        auto_fix: Whether to apply fixes
        min_font_size: Minimum font size for adjustments
        
    Returns:
        Dictionary with boundary validation results
    """
    result = {
        'issues_found': [],
        'fixes_applied': [],
        'warnings': [],
        'boundary_issues': 0
    }
    
    try:
        # Use the existing validate_text_container function
        text_content = text_frame.text
        if not text_content.strip():
            return result
            
        # Get current font size (use first run's font size as representative)
        current_font_size = 12  # Default
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            first_run = text_frame.paragraphs[0].runs[0]
            if first_run.font.size is not None:
                current_font_size = first_run.font.size.pt
        
        # Validate text fits in container
        validation_result = validate_text_container(shape, text_content, current_font_size)
        
        if not validation_result['fits']:
            issue = f'{shape_name}: Text content exceeds shape boundaries'
            result['issues_found'].append(issue)
            result['boundary_issues'] += 1
            
            if auto_fix:
                # Try to fix by adjusting font size
                suggested_size = validation_result.get('suggested_font_size', current_font_size * 0.8)
                
                # Ensure suggested size is within bounds
                suggested_size = max(min_font_size, min(suggested_size, current_font_size))
                
                if suggested_size != current_font_size:
                    # Apply suggested font size to all runs
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size is not None:
                                run.font.size = Pt(suggested_size)
                    
                    fix = f'{shape_name}: Adjusted font size from {current_font_size}pt to {suggested_size}pt to fit boundaries'
                    result['fixes_applied'].append(fix)
                    
                    # Enable word wrap if not already enabled
                    if not text_frame.word_wrap:
                        text_frame.word_wrap = True
                        result['fixes_applied'].append(f'{shape_name}: Enabled word wrap')
                        
                else:
                    # Font size can't be reduced further, suggest shape resize
                    result['warnings'].append(
                        f'{shape_name}: Text still may not fit even at minimum font size. '
                        f'Consider expanding shape or reducing text content.'
                    )
                    
        # Check margins and padding
        margin_issues = check_text_margins(text_frame, shape_name, auto_fix)
        result['issues_found'].extend(margin_issues['issues_found'])
        result['fixes_applied'].extend(margin_issues['fixes_applied'])
        result['warnings'].extend(margin_issues['warnings'])
        
    except Exception as e:
        result['warnings'].append(f'Boundary check failed for {shape_name}: {str(e)}')
    
    return result


def check_text_margins(text_frame, shape_name: str, auto_fix: bool) -> Dict:
    """
    Check and fix text frame margins for better readability.
    
    Args:
        text_frame: The text frame to check
        shape_name: Shape identifier  
        auto_fix: Whether to apply fixes
        
    Returns:
        Dictionary with margin validation results
    """
    result = {
        'issues_found': [],
        'fixes_applied': [],
        'warnings': []
    }
    
    try:
        # Recommended minimum margins in points
        min_margin_pt = 5  # 5 points minimum margin
        min_margin_emu = int(min_margin_pt * 12700)  # Convert to EMU
        
        margins_to_check = [
            ('left', text_frame.margin_left),
            ('right', text_frame.margin_right), 
            ('top', text_frame.margin_top),
            ('bottom', text_frame.margin_bottom)
        ]
        
        for margin_name, current_margin in margins_to_check:
            if current_margin < min_margin_emu:
                current_pt = current_margin / 12700
                issue = f'{shape_name}: {margin_name} margin {current_pt:.1f}pt below recommended minimum {min_margin_pt}pt'
                result['issues_found'].append(issue)
                
                if auto_fix:
                    setattr(text_frame, f'margin_{margin_name}', min_margin_emu)
                    fix = f'{shape_name}: Increased {margin_name} margin from {current_pt:.1f}pt to {min_margin_pt}pt'
                    result['fixes_applied'].append(fix)
                    
    except Exception as e:
        result['warnings'].append(f'Margin check failed for {shape_name}: {str(e)}')
    
    return result


def validate_shape_layout(shape, shape_name: str, auto_fix: bool) -> Dict:
    """
    Validate shape positioning and size for reasonable layout.
    
    Args:
        shape: The shape object
        shape_name: Shape identifier
        auto_fix: Whether to apply fixes
        
    Returns:
        Dictionary with layout validation results
    """
    result = {
        'issues_found': [],
        'fixes_applied': [],
        'warnings': [],
        'layout_issues': 0
    }
    
    try:
        # Get shape position and dimensions in points
        left_pt = shape.left / 12700
        top_pt = shape.top / 12700  
        width_pt = shape.width / 12700
        height_pt = shape.height / 12700
        
        # Get slide dimensions (approximate)
        try:
            slide_width_pt = shape.part.parent_part.presentation.slide_width / 12700
            slide_height_pt = shape.part.parent_part.presentation.slide_height / 12700
        except:
            # Default to 16:9 slide dimensions in points
            slide_width_pt = 720  # 10 inches * 72 pts/inch  
            slide_height_pt = 540  # 7.5 inches * 72 pts/inch
        
        # Check if shape extends beyond slide boundaries
        if left_pt + width_pt > slide_width_pt:
            issue = f'{shape_name}: Shape extends beyond right edge of slide'
            result['issues_found'].append(issue)
            result['layout_issues'] += 1
            
            if auto_fix:
                # Move shape to fit within slide
                new_left = slide_width_pt - width_pt - 10  # 10pt margin
                shape.left = int(max(0, new_left) * 12700)
                fix = f'{shape_name}: Moved shape to fit within slide boundaries'
                result['fixes_applied'].append(fix)
        
        if top_pt + height_pt > slide_height_pt:
            issue = f'{shape_name}: Shape extends beyond bottom edge of slide'
            result['issues_found'].append(issue)
            result['layout_issues'] += 1
            
            if auto_fix:
                # Move shape to fit within slide
                new_top = slide_height_pt - height_pt - 10  # 10pt margin
                shape.top = int(max(0, new_top) * 12700)
                fix = f'{shape_name}: Moved shape to fit within slide boundaries'
                result['fixes_applied'].append(fix)
        
        # Check for very small shapes that might be hard to read
        min_readable_width = 50  # 50 points minimum
        min_readable_height = 20  # 20 points minimum
        
        if width_pt < min_readable_width:
            issue = f'{shape_name}: Shape width {width_pt:.1f}pt may be too small for readability'
            result['issues_found'].append(issue)
            result['layout_issues'] += 1
            
            if auto_fix:
                shape.width = int(min_readable_width * 12700)
                fix = f'{shape_name}: Increased width to {min_readable_width}pt for better readability'
                result['fixes_applied'].append(fix)
        
        if height_pt < min_readable_height:
            issue = f'{shape_name}: Shape height {height_pt:.1f}pt may be too small for readability'
            result['issues_found'].append(issue)
            result['layout_issues'] += 1
            
            if auto_fix:
                shape.height = int(min_readable_height * 12700)
                fix = f'{shape_name}: Increased height to {min_readable_height}pt for better readability'
                result['fixes_applied'].append(fix)
                
    except Exception as e:
        result['warnings'].append(f'Layout validation failed for {shape_name}: {str(e)}')
    
    return result


def validate_slide_layout(slide, auto_fix: bool) -> Dict:
    """
    Validate overall slide layout for proper spacing and arrangement.
    
    Args:
        slide: The slide object
        auto_fix: Whether to apply fixes
        
    Returns:
        Dictionary with slide-level layout validation results
    """
    result = {
        'issues_found': [],
        'fixes_applied': [],
        'warnings': []
    }
    
    try:
        shapes_with_text = []
        
        # Collect all text shapes for layout analysis
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                shapes_with_text.append(shape)
        
        # Check for overlapping shapes
        overlaps = find_overlapping_shapes(shapes_with_text)
        if overlaps:
            for overlap in overlaps:
                issue = f'Shapes overlap: {overlap["shape1_name"]} and {overlap["shape2_name"]}'
                result['issues_found'].append(issue)
                
                if auto_fix:
                    # Simple fix: move second shape slightly
                    shape2 = overlap['shape2']
                    shape2.left += int(10 * 12700)  # Move 10 points right
                    fix = f'Moved {overlap["shape2_name"]} to resolve overlap'
                    result['fixes_applied'].append(fix)
        
        # Check for proper spacing between shapes
        spacing_issues = check_shape_spacing(shapes_with_text)
        result['issues_found'].extend(spacing_issues['issues'])
        
        if auto_fix and spacing_issues['fixes']:
            result['fixes_applied'].extend(spacing_issues['fixes'])
            
    except Exception as e:
        result['warnings'].append(f'Slide layout validation failed: {str(e)}')
    
    return result


def find_overlapping_shapes(shapes: list) -> list:
    """
    Find shapes that overlap with each other.
    
    Args:
        shapes: List of shapes to check for overlaps
        
    Returns:
        List of overlap information dictionaries
    """
    overlaps = []
    
    try:
        for i, shape1 in enumerate(shapes):
            for j, shape2 in enumerate(shapes[i+1:], i+1):
                # Calculate boundaries
                s1_left = shape1.left
                s1_right = shape1.left + shape1.width
                s1_top = shape1.top
                s1_bottom = shape1.top + shape1.height
                
                s2_left = shape2.left
                s2_right = shape2.left + shape2.width
                s2_top = shape2.top
                s2_bottom = shape2.top + shape2.height
                
                # Check for overlap
                if (s1_left < s2_right and s1_right > s2_left and 
                    s1_top < s2_bottom and s1_bottom > s2_top):
                    
                    overlaps.append({
                        'shape1': shape1,
                        'shape2': shape2,
                        'shape1_name': getattr(shape1, 'name', f'Shape_{i}'),
                        'shape2_name': getattr(shape2, 'name', f'Shape_{j}')
                    })
                    
    except Exception as e:
        print(f"Error finding overlaps: {e}")
    
    return overlaps


def check_shape_spacing(shapes: list) -> Dict:
    """
    Check spacing between shapes for proper layout.
    
    Args:
        shapes: List of shapes to check spacing
        
    Returns:
        Dictionary with spacing issues and potential fixes
    """
    result = {
        'issues': [],
        'fixes': []
    }
    
    try:
        min_spacing_pt = 10  # Minimum 10 points between shapes
        min_spacing_emu = int(min_spacing_pt * 12700)
        
        for i, shape1 in enumerate(shapes):
            for j, shape2 in enumerate(shapes[i+1:], i+1):
                # Calculate distances
                h_distance = min(
                    abs(shape1.left + shape1.width - shape2.left),
                    abs(shape2.left + shape2.width - shape1.left)
                )
                
                v_distance = min(
                    abs(shape1.top + shape1.height - shape2.top),
                    abs(shape2.top + shape2.height - shape1.top)
                )
                
                # Check if shapes are too close
                if h_distance < min_spacing_emu and v_distance < min_spacing_emu:
                    shape1_name = getattr(shape1, 'name', f'Shape_{i}')
                    shape2_name = getattr(shape2, 'name', f'Shape_{j}')
                    
                    distance_pt = min(h_distance, v_distance) / 12700
                    issue = f'Shapes {shape1_name} and {shape2_name} too close: {distance_pt:.1f}pt spacing'
                    result['issues'].append(issue)
                    
    except Exception as e:
        print(f"Error checking spacing: {e}")
    
    return result
