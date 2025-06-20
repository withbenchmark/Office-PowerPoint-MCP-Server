#!/usr/bin/env python
"""
MCP Server for PowerPoint manipulation using python-pptx.
"""
import os
import json
import tempfile
from typing import Dict, List, Optional, Any, Union
from mcp.server.fastmcp import FastMCP
import argparse

import ppt_utils

# Initialize the FastMCP server
app = FastMCP(
    name="ppt-mcp-server",
    description="MCP Server for PowerPoint manipulation using python-pptx",
    version="1.0.0", 
    log_level="INFO"
)

# Global state to store presentations in memory
presentations = {}
current_presentation_id = None

# Template configuration
def get_template_search_directories():
    """
    Get list of directories to search for templates.
    Uses environment variable PPT_TEMPLATE_PATH if set, otherwise uses default directories.
    
    Returns:
        List of directories to search for templates
    """
    template_env_path = os.environ.get('PPT_TEMPLATE_PATH')
    
    if template_env_path:
        # If environment variable is set, use it as the primary template directory
        # Support multiple paths separated by colon (Unix) or semicolon (Windows)
        import platform
        separator = ';' if platform.system() == "Windows" else ':'
        env_dirs = [path.strip() for path in template_env_path.split(separator) if path.strip()]
        
        # Verify that the directories exist
        valid_env_dirs = []
        for dir_path in env_dirs:
            expanded_path = os.path.expanduser(dir_path)
            if os.path.exists(expanded_path) and os.path.isdir(expanded_path):
                valid_env_dirs.append(expanded_path)
        
        if valid_env_dirs:
            # Add default fallback directories
            return valid_env_dirs + ['.', './templates', './assets', './resources']
        else:
            print(f"Warning: PPT_TEMPLATE_PATH directories not found: {template_env_path}")
    
    # Default search directories when no environment variable or invalid paths
    return ['.', './templates', './assets', './resources']

# ---- Helper Functions ----

def get_current_presentation():
    """Get the current presentation object or raise an error if none is loaded."""
    if current_presentation_id is None or current_presentation_id not in presentations:
        raise ValueError("No presentation is currently loaded. Please create or open a presentation first.")
    return presentations[current_presentation_id]

def validate_parameters(params):
    """
    Validate parameters against constraints.
    
    Args:
        params: Dictionary of parameter name: (value, constraints) pairs
        
    Returns:
        (True, None) if all valid, or (False, error_message) if invalid
    """
    for param_name, (value, constraints) in params.items():
        for constraint_func, error_msg in constraints:
            if not constraint_func(value):
                return False, f"Parameter '{param_name}': {error_msg}"
    return True, None

def is_positive(value):
    """Check if a value is positive."""
    return value > 0

def is_non_negative(value):
    """Check if a value is non-negative."""
    return value >= 0

def is_in_range(min_val, max_val):
    """Create a function that checks if a value is in a range."""
    return lambda x: min_val <= x <= max_val

def is_in_list(valid_list):
    """Create a function that checks if a value is in a list."""
    return lambda x: x in valid_list

def is_valid_rgb(color_list):
    """Check if a color list is a valid RGB tuple."""
    if not isinstance(color_list, list) or len(color_list) != 3:
        return False
    return all(isinstance(c, int) and 0 <= c <= 255 for c in color_list)

def add_shape_direct(slide, shape_type: str, left: float, top: float, width: float, height: float) -> Any:
    """
    Add an auto shape to a slide using direct integer values instead of enum objects.
    
    This implementation provides a reliable alternative that bypasses potential 
    enum-related issues in the python-pptx library.
    
    Args:
        slide: The slide object
        shape_type: Shape type string (e.g., 'rectangle', 'oval', 'triangle')
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created shape
    """
    from pptx.util import Inches
    
    # Direct mapping of shape types to their integer values
    # These values are directly from the MS Office VBA documentation
    shape_type_map = {
        'rectangle': 1,
        'rounded_rectangle': 2, 
        'oval': 9,
        'diamond': 4,
        'triangle': 5,  # This is ISOSCELES_TRIANGLE
        'right_triangle': 6,
        'pentagon': 56,
        'hexagon': 10,
        'heptagon': 11,
        'octagon': 12,
        'star': 12,  # This is STAR_5_POINTS (value 12)
        'arrow': 13,
        'cloud': 35,
        'heart': 21,
        'lightning_bolt': 22,
        'sun': 23,
        'moon': 24,
        'smiley_face': 17,
        'no_symbol': 19,
        'flowchart_process': 112,
        'flowchart_decision': 114,
        'flowchart_data': 115,
        'flowchart_document': 119
    }
    
    # Check if shape type is valid before trying to use it
    shape_type_lower = str(shape_type).lower()
    if shape_type_lower not in shape_type_map:
        available_shapes = ', '.join(sorted(shape_type_map.keys()))
        raise ValueError(f"Unsupported shape type: '{shape_type}'. Available shape types: {available_shapes}")
    
    # Get the integer value for the shape type
    shape_value = shape_type_map[shape_type_lower]
    
    # Create the shape using the direct integer value
    try:
        # The integer value is passed directly to add_shape
        shape = slide.shapes.add_shape(
            shape_value, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        return shape
    except Exception as e:
        raise ValueError(f"Failed to create '{shape_type}' shape using direct value {shape_value}: {str(e)}")

# ---- Presentation Tools ----

@app.tool()
def create_presentation(id: Optional[str] = None) -> Dict:
    """Create a new PowerPoint presentation."""
    global current_presentation_id
    
    # Create a new presentation
    pres = ppt_utils.create_presentation()
    
    # Generate an ID if not provided
    if id is None:
        id = f"presentation_{len(presentations) + 1}"
    
    # Store the presentation
    presentations[id] = pres
    current_presentation_id = id
    
    return {
        "presentation_id": id,
        "message": f"Created new presentation with ID: {id}",
        "slide_count": len(pres.slides)
    }

@app.tool()
def create_presentation_from_template(template_path: str, id: Optional[str] = None) -> Dict:
    """Create a new PowerPoint presentation from a template file."""
    global current_presentation_id
    
    # Check if template file exists
    if not os.path.exists(template_path):
        # Try to find the template by searching in configured directories
        search_dirs = get_template_search_directories()
        template_name = os.path.basename(template_path)
        
        for directory in search_dirs:
            potential_path = os.path.join(directory, template_name)
            if os.path.exists(potential_path):
                template_path = potential_path
                break
        else:
            env_path_info = f" (PPT_TEMPLATE_PATH: {os.environ.get('PPT_TEMPLATE_PATH', 'not set')})" if os.environ.get('PPT_TEMPLATE_PATH') else ""
            return {
                "error": f"Template file not found: {template_path}. Searched in {', '.join(search_dirs)}{env_path_info}"
            }
    
    # Create presentation from template
    try:
        pres = ppt_utils.create_presentation_from_template(template_path)
    except Exception as e:
        return {
            "error": f"Failed to create presentation from template: {str(e)}"
        }
    
    # Generate an ID if not provided
    if id is None:
        id = f"presentation_{len(presentations) + 1}"
    
    # Store the presentation
    presentations[id] = pres
    current_presentation_id = id
    
    return {
        "presentation_id": id,
        "message": f"Created new presentation from template '{template_path}' with ID: {id}",
        "template_path": template_path,
        "slide_count": len(pres.slides),
        "layout_count": len(pres.slide_layouts)
    }

@app.tool()
def open_presentation(file_path: str, id: Optional[str] = None) -> Dict:
    """Open an existing PowerPoint presentation from a file."""
    global current_presentation_id
    
    # Check if file exists
    if not os.path.exists(file_path):
        return {
            "error": f"File not found: {file_path}"
        }
    
    # Open the presentation
    try:
        pres = ppt_utils.open_presentation(file_path)
    except Exception as e:
        return {
            "error": f"Failed to open presentation: {str(e)}"
        }
    
    # Generate an ID if not provided
    if id is None:
        id = f"presentation_{len(presentations) + 1}"
    
    # Store the presentation
    presentations[id] = pres
    current_presentation_id = id
    
    return {
        "presentation_id": id,
        "message": f"Opened presentation from {file_path} with ID: {id}",
        "slide_count": len(pres.slides)
    }

@app.tool()
def save_presentation(file_path: str, presentation_id: Optional[str] = None) -> Dict:
    """Save a presentation to a file."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    # Save the presentation
    try:
        saved_path = ppt_utils.save_presentation(presentations[pres_id], file_path)
        return {
            "message": f"Presentation saved to {saved_path}",
            "file_path": saved_path
        }
    except Exception as e:
        return {
            "error": f"Failed to save presentation: {str(e)}"
        }

@app.tool()
def get_presentation_info(presentation_id: Optional[str] = None) -> Dict:
    """Get information about a presentation."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Get slide layouts
    layouts = ppt_utils.get_slide_layouts(pres)
    
    # Get core properties
    core_props = ppt_utils.get_core_properties(pres)
    
    return {
        "presentation_id": pres_id,
        "slide_count": len(pres.slides),
        "slide_layouts": layouts,
        "core_properties": core_props
    }

@app.tool()
def set_core_properties(
    title: Optional[str] = None,
    subject: Optional[str] = None,
    author: Optional[str] = None,
    keywords: Optional[str] = None,
    comments: Optional[str] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Set core document properties."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Set core properties
    try:
        ppt_utils.set_core_properties(
            pres, title=title, subject=subject, author=author, 
            keywords=keywords, comments=comments
        )
        
        # Get updated properties
        updated_props = ppt_utils.get_core_properties(pres)
        
        return {
            "message": "Core properties updated successfully",
            "core_properties": updated_props
        }
    except Exception as e:
        return {
            "error": f"Failed to set core properties: {str(e)}"
        }

@app.tool() 
def get_template_info(template_path: str) -> Dict:
    """Get information about a template file including layouts and properties."""
    # Check if template file exists
    if not os.path.exists(template_path):
        # Try to find the template by searching in configured directories
        search_dirs = get_template_search_directories()
        template_name = os.path.basename(template_path)
        
        for directory in search_dirs:
            potential_path = os.path.join(directory, template_name)
            if os.path.exists(potential_path):
                template_path = potential_path
                break
        else:
            env_path_info = f" (PPT_TEMPLATE_PATH: {os.environ.get('PPT_TEMPLATE_PATH', 'not set')})" if os.environ.get('PPT_TEMPLATE_PATH') else ""
            return {
                "error": f"Template file not found: {template_path}. Searched in {', '.join(search_dirs)}{env_path_info}"
            }
    
    # Get template information
    try:
        template_info = ppt_utils.get_template_info(template_path)
        return template_info
    except Exception as e:
        return {
            "error": f"Failed to read template info: {str(e)}"
        }

# ---- Slide Tools ----

@app.tool()
def add_slide(
    layout_index: int = 1,
    title: Optional[str] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Add a new slide to the presentation."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate layout index
    if layout_index < 0 or layout_index >= len(pres.slide_layouts):
        return {
            "error": f"Invalid layout index: {layout_index}. Available layouts: 0-{len(pres.slide_layouts) - 1}",
            "available_layouts": ppt_utils.get_slide_layouts(pres)
        }
    
    # Add the slide
    slide, error = ppt_utils.safe_operation(
        "add_slide",
        lambda: ppt_utils.add_slide(pres, layout_index)
    )
    
    if error:
        return {"error": error}
    
    # Set the title if provided
    if title and slide[0].shapes.title:
        _, error = ppt_utils.safe_operation(
            "set_title",
            lambda: ppt_utils.set_title(slide[0], title)
        )
        if error:
            return {
                "warning": f"Slide created but failed to set title: {error}",
                "slide_index": len(pres.slides) - 1,
                "layout_name": slide[1].name
            }
    
    # Get placeholders
    placeholders, error = ppt_utils.safe_operation(
        "get_placeholders",
        lambda: ppt_utils.get_placeholders(slide[0])
    )
    
    if error:
        placeholders = []
    
    # Automatically validate and fix the slide content
    validation_result, error = ppt_utils.safe_operation(
        "validate_slide",
        lambda: ppt_utils.validate_and_fix_slide_content(slide[0], auto_fix=True)
    )
    
    result = {
        "message": f"Added slide with layout '{slide[1].name}'",
        "slide_index": len(pres.slides) - 1,
        "layout_name": slide[1].name,
        "placeholders": placeholders
    }
    
    # Include validation results if available
    if validation_result and not error:
        result["validation"] = {
            "issues_found": len(validation_result.get('issues_found', [])),
            "fixes_applied": len(validation_result.get('fixes_applied', [])),
            "summary": validation_result.get('validation_summary', {}),
            "message": validation_result.get('message', 'Validation completed')
        }
        
        # Add warnings if any fixes were applied
        if validation_result.get('fixes_applied'):
            result["validation"]["fixes_applied_details"] = validation_result['fixes_applied']
            
    elif error:
        result["validation"] = {"warning": f"Slide validation failed: {error}"}
    
    return result

@app.tool()
def get_slide_info(slide_index: int, presentation_id: Optional[str] = None) -> Dict:
    """Get information about a specific slide."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Get placeholders
    placeholders = ppt_utils.get_placeholders(slide)
    
    # Get shapes information
    shapes_info = []
    for i, shape in enumerate(slide.shapes):
        shape_info = {
            "index": i,
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "width": shape.width.inches,
            "height": shape.height.inches,
            "left": shape.left.inches,
            "top": shape.top.inches
        }
        shapes_info.append(shape_info)
    
    return {
        "slide_index": slide_index,
        "placeholders": placeholders,
        "shapes": shapes_info
    }

@app.tool()
def populate_placeholder(
    slide_index: int,
    placeholder_idx: int,
    text: str,
    presentation_id: Optional[str] = None
) -> Dict:
    """Populate a placeholder with text."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    try:
        # Check if placeholder exists
        if placeholder_idx not in [p.placeholder_format.idx for p in slide.placeholders]:
            return {
                "error": f"Placeholder with index {placeholder_idx} not found in slide {slide_index}"
            }
        
        # Populate the placeholder
        ppt_utils.populate_placeholder(slide, placeholder_idx, text)
        
        return {
            "message": f"Populated placeholder {placeholder_idx} in slide {slide_index}"
        }
    except Exception as e:
        return {
            "error": f"Failed to populate placeholder: {str(e)}"
        }

@app.tool()
def add_bullet_points(
    slide_index: int,
    placeholder_idx: int,
    bullet_points: List[str],
    presentation_id: Optional[str] = None
) -> Dict:
    """Add bullet points to a placeholder."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    try:
        # Check if placeholder exists
        if placeholder_idx not in [p.placeholder_format.idx for p in slide.placeholders]:
            return {
                "error": f"Placeholder with index {placeholder_idx} not found in slide {slide_index}"
            }
        
        # Get the placeholder
        placeholder = slide.placeholders[placeholder_idx]
        
        # Add bullet points
        ppt_utils.add_bullet_points(placeholder, bullet_points)
        
        return {
            "message": f"Added {len(bullet_points)} bullet points to placeholder {placeholder_idx} in slide {slide_index}"
        }
    except Exception as e:
        return {
            "error": f"Failed to add bullet points: {str(e)}"
        }

# ---- Text Tools ----

@app.tool()
def add_textbox(
    slide_index: int,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: Optional[int] = None,
    font_name: Optional[str] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    color: Optional[List[int]] = None,
    alignment: Optional[str] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Add a textbox to a slide."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    try:
        # Add the textbox
        textbox = ppt_utils.add_textbox(slide, left, top, width, height, text)
        
        # Format the text if formatting options are provided
        if any([font_size, font_name, bold, italic, color, alignment]):
            ppt_utils.format_text(
                textbox.text_frame,
                font_size=font_size,
                font_name=font_name,
                bold=bold,
                italic=italic,
                color=tuple(color) if color else None,
                alignment=alignment
            )
        
        # Automatically validate and fix the slide content after adding textbox
        validation_result, error = ppt_utils.safe_operation(
            "validate_slide",
            lambda: ppt_utils.validate_and_fix_slide_content(slide, auto_fix=True)
        )
        
        result = {
            "message": f"Added textbox to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1
        }
        
        # Include validation results if available
        if validation_result and not error:
            if validation_result.get('fixes_applied'):
                result["validation"] = {
                    "fixes_applied": len(validation_result['fixes_applied']),
                    "message": "Textbox automatically adjusted for optimal layout",
                    "details": validation_result['fixes_applied']
                }
        elif error:
            result["validation_warning"] = f"Post-creation validation failed: {error}"
            
        return result
    except Exception as e:
        return {
            "error": f"Failed to add textbox: {str(e)}"
        }

# ---- Image Tools ----

@app.tool()
def add_image(
    slide_index: int,
    image_path: str,
    left: float,
    top: float,
    width: Optional[float] = None,
    height: Optional[float] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Add an image to a slide with graceful error recovery."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Check if image file exists
    if not os.path.exists(image_path):
        # Try to find the image by searching in common directories
        common_dirs = ['.', './images', './assets', './resources']
        image_name = os.path.basename(image_path)
        
        for directory in common_dirs:
            potential_path = os.path.join(directory, image_name)
            if os.path.exists(potential_path):
                image_path = potential_path
                break
        else:
            return {
                "error": f"Image file not found: {image_path}. Searched in {', '.join(common_dirs)}"
            }
    
    # Define multiple approaches to add the image
    def add_with_size():
        return ppt_utils.add_image(slide, image_path, left, top, width, height)
        
    def add_without_size():
        return ppt_utils.add_image(slide, image_path, left, top)
    
    def add_with_pil():
        from PIL import Image
        img = Image.open(image_path)
        img_width, img_height = img.size
        
        # Calculate aspect ratio and use it to determine missing dimension
        aspect_ratio = img_width / img_height
        
        if width is not None and height is None:
            h = width / aspect_ratio
            return ppt_utils.add_image(slide, image_path, left, top, width, h)
        elif height is not None and width is None:
            w = height * aspect_ratio
            return ppt_utils.add_image(slide, image_path, left, top, w, height)
        else:
            return ppt_utils.add_image(slide, image_path, left, top, width, height)
    
    approaches = [
        (add_with_size, "Adding image with specified dimensions"),
        (add_without_size, "Adding image with original dimensions"),
        (add_with_pil, "Adding image with calculated dimensions using PIL")
    ]
    
    picture, error = ppt_utils.try_multiple_approaches("add image", approaches)
    
    if error:
        return {
            "error": error
        }
    
    return {
        "message": f"Added image to slide {slide_index}",
        "shape_index": len(slide.shapes) - 1,
        "width": picture.width.inches,
        "height": picture.height.inches
    }

@app.tool()
def add_image_from_base64(
    slide_index: int,
    base64_string: str,
    left: float,
    top: float,
    width: Optional[float] = None,
    height: Optional[float] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Add an image from a base64 encoded string to a slide."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    try:
        # Add the image
        picture = ppt_utils.add_image_from_base64(slide, base64_string, left, top, width, height)
        
        return {
            "message": f"Added image to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1,
            "width": picture.width.inches,
            "height": picture.height.inches
        }
    except Exception as e:
        return {
            "error": f"Failed to add image: {str(e)}"
        }

# ---- Table Tools ----

@app.tool()
def add_table(
    slide_index: int,
    rows: int,
    cols: int,
    left: float,
    top: float,
    width: float,
    height: float,
    data: Optional[List[List[str]]] = None,
    presentation_id: Optional[str] = None,
    filename: Optional[str] = None  # Added to satisfy MCP framework validation
) -> Dict:
    """Add a table to a slide with comprehensive parameter validation."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate parameters
    valid, error = validate_parameters({
        "rows": (rows, [(is_positive, "must be a positive integer")]),
        "cols": (cols, [(is_positive, "must be a positive integer")]),
        "left": (left, [(is_non_negative, "must be non-negative")]),
        "top": (top, [(is_non_negative, "must be non-negative")]),
        "width": (width, [(is_positive, "must be positive")]),
        "height": (height, [(is_positive, "must be positive")]),
    })
    
    if not valid:
        return {"error": error}
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate data if provided
    if data is not None:
        if not isinstance(data, list):
            return {"error": "Data must be a list of rows"}
        
        for i, row in enumerate(data):
            if not isinstance(row, list):
                return {"error": f"Row {i} must be a list of cell values"}
    
    try:
        # Add the table
        table, error = ppt_utils.safe_operation(
            "add_table",
            lambda: ppt_utils.add_table(slide, rows, cols, left, top, width, height)
        )
        
        if error:
            return {"error": error}
        
        # Populate the table if data is provided
        warnings = []
        if data:
            for row_idx, row_data in enumerate(data):
                if row_idx >= rows:
                    warnings.append(f"Ignored excess data: table has only {rows} rows but data has {len(data)} rows")
                    break
                    
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx >= cols:
                        warnings.append(f"Ignored excess data in row {row_idx}: table has only {cols} columns")
                        break
                        
                    _, cell_error = ppt_utils.safe_operation(
                        f"set_cell_text(row={row_idx}, col={col_idx})",
                        lambda: ppt_utils.set_cell_text(table, row_idx, col_idx, str(cell_text))
                    )
                    
                    if cell_error:
                        warnings.append(cell_error)
        
        result = {
            "message": f"Added {rows}x{cols} table to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1
        }
        
        if warnings:
            result["warnings"] = warnings
            
        return result
    except Exception as e:
        return {
            "error": f"Failed to add table: {str(e)}"
        }

@app.tool()
def format_table_cell(
    slide_index: int,
    shape_index: int,
    row: int,
    col: int,
    font_size: Optional[int] = None,
    font_name: Optional[str] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    color: Optional[List[int]] = None,
    bg_color: Optional[List[int]] = None,
    alignment: Optional[str] = None,
    vertical_alignment: Optional[str] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Format a table cell with comprehensive error handling and parameter validation.
    
    This function applies formatting to a cell in a table on a slide. It provides options
    for text formatting (font size, name, style, color), cell background color, and
    text alignment.
    
    Args:
        slide_index: Index of the slide containing the table (0-based)
        shape_index: Index of the table shape on the slide (0-based)
        row: Row index of the cell to format (0-based)
        col: Column index of the cell to format (0-based)
        font_size: Font size in points
        font_name: Font name/family (e.g., 'Arial', 'Calibri')
        bold: Whether text should be bold (True/False)
        italic: Whether text should be italic (True/False)
        color: RGB color list for text [R, G, B] (0-255 for each value)
        bg_color: RGB color list for cell background [R, G, B] (0-255 for each value)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical text alignment ('top', 'middle', 'bottom')
        presentation_id: ID of the presentation to use (uses current presentation if not specified)
    
    Returns:
        Dict with keys:
            - message: Success message
            - error: Error message if operation failed
            - warning: Warning message if some formatting was applied but some failed
    
    Examples:
        To format a header cell with bold text and gray background:
            format_table_cell(0, 1, 0, 1, font_size=14, bold=True, bg_color=[200, 200, 200])
            
        To center text in a cell:
            format_table_cell(0, 1, 2, 3, alignment="center", vertical_alignment="middle")
    """
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Check if shape index is valid
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Validate parameters
    valid_alignments = ['left', 'center', 'right', 'justify']
    valid_vertical_alignments = ['top', 'middle', 'bottom']
    
    validations = {}
    
    if font_size is not None:
        validations["font_size"] = (font_size, [(is_positive, "must be a positive integer")])
    
    if alignment is not None:
        validations["alignment"] = (alignment.lower(), [(lambda x: x in valid_alignments, 
                                    f"must be one of {', '.join(valid_alignments)}")])
    
    if vertical_alignment is not None:
        validations["vertical_alignment"] = (vertical_alignment.lower(), 
                                           [(lambda x: x in valid_vertical_alignments, 
                                             f"must be one of {', '.join(valid_vertical_alignments)}")])
    
    if color is not None:
        validations["color"] = (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
    
    if bg_color is not None:
        validations["bg_color"] = (bg_color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
    
    if validations:
        valid, error = validate_parameters(validations)
        if not valid:
            return {"error": error}
    
    try:
        # Check if shape is a table
        if not hasattr(shape, 'table'):
            # Try to recover if this is a graphic frame containing a table
            if hasattr(shape, 'graphic') and hasattr(shape.graphic, 'graphicData'):
                # This might be a table in a graphic frame
                warnings = ["Shape is not directly a table, attempting to extract table from graphic frame"]
                # Further recovery logic would be needed here
                return {
                    "error": "Shape at index {shape_index} is not a table",
                    "warning": "If this is a table, it might be in a graphic frame which requires different handling"
                }
            else:
                return {
                    "error": f"Shape at index {shape_index} is not a table"
                }
        
        table = shape.table
        
        # Check if row and column indices are valid
        if row < 0 or row >= len(table.rows):
            return {
                "error": f"Invalid row index: {row}. Available rows: 0-{len(table.rows) - 1}"
            }
            
        if col < 0 or col >= len(table.columns):
            return {
                "error": f"Invalid column index: {col}. Available columns: 0-{len(table.columns) - 1}"
            }
        
        # Get the cell
        cell = table.cell(row, col)
        
        # Format the cell with error handling
        warnings = []
        
        # Try multiple formatting operations and collect any warnings
        try:
            ppt_utils.format_table_cell(
                cell,
                font_size=font_size,
                font_name=font_name,
                bold=bold,
                italic=italic,
                color=tuple(color) if color else None,
                bg_color=tuple(bg_color) if bg_color else None,
                alignment=alignment,
                vertical_alignment=vertical_alignment
            )
        except Exception as e:
            # Try individual formatting operations to recover
            formatting_ops = [
                (lambda: ppt_utils.format_text(cell.text_frame, font_size=font_size, font_name=font_name, 
                                              bold=bold, italic=italic, 
                                              color=tuple(color) if color else None,
                                              alignment=alignment),
                 "text formatting"),
                
                (lambda: cell.fill.solid() if bg_color else None, "background preparation"),
                
                (lambda: setattr(cell.fill.fore_color, 'rgb', 
                                tuple(bg_color)) if bg_color else None, 
                 "background color"),
                
                (lambda: setattr(cell.text_frame, 'vertical_anchor', 
                                ppt_utils.vertical_alignment_map.get(vertical_alignment)) 
                 if vertical_alignment else None,
                 "vertical alignment")
            ]
            
            for op_func, op_name in formatting_ops:
                try:
                    op_func()
                except Exception as sub_e:
                    warnings.append(f"Failed to apply {op_name}: {str(sub_e)}")
        
        result = {
            "message": f"Formatted cell at row {row}, column {col} in table at shape index {shape_index} on slide {slide_index}"
        }
        
        if warnings:
            result["warnings"] = warnings
            
        return result
    except Exception as e:
        return {
            "error": f"Failed to format table cell: {str(e)}"
        }

# ---- Shape Tools ----

@app.tool()
def add_shape(
    slide_index: int,
    shape_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: Optional[List[int]] = None,
    line_color: Optional[List[int]] = None,
    line_width: Optional[float] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Add an auto shape to a slide.
    
    This function adds a shape to a slide in the presentation. It supports various shape types
    and allows customization of fill color, line color, and line width.
    
    Args:
        slide_index: Index of the slide to add the shape to (0-based)
        shape_type: Type of shape to add. Supported types include:
            - Basic shapes: 'rectangle', 'rounded_rectangle', 'oval', 'triangle', 'diamond'
            - Polygons: 'pentagon', 'hexagon', 'heptagon', 'octagon'
            - Stars and arrows: 'star', 'arrow'
            - Misc: 'cloud', 'heart', 'lightning_bolt', 'sun', 'moon', 'smiley_face', 'no_symbol'
            - Flowchart: 'flowchart_process', 'flowchart_decision', 'flowchart_data'
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        fill_color: RGB color list for shape fill [R, G, B] (0-255 for each value)
        line_color: RGB color list for shape outline [R, G, B] (0-255 for each value)
        line_width: Width of the shape outline in points
        presentation_id: ID of the presentation to use (uses current presentation if not specified)
    """
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    try:
        # Use the direct implementation that bypasses the enum issues
        shape = add_shape_direct(slide, shape_type, left, top, width, height)
        
        # Format the shape if formatting options are provided
        if any([fill_color, line_color, line_width]):
            ppt_utils.format_shape(
                shape,
                fill_color=tuple(fill_color) if fill_color else None,
                line_color=tuple(line_color) if line_color else None,
                line_width=line_width
            )
        
        return {
            "message": f"Added {shape_type} shape to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1
        }
    except ValueError as e:
        # Specific handling for validation errors
        return {
            "error": str(e)
        }
    except Exception as e:
        return {
            "error": f"Failed to add shape '{shape_type}': {str(e)}"
        }

# ---- Chart Tools ----

@app.tool()
def add_chart(
    slide_index: int,
    chart_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    categories: List[str],
    series_names: List[str],
    series_values: List[List[float]],
    has_legend: bool = True,
    legend_position: str = "right",
    has_data_labels: bool = False,
    title: Optional[str] = None,
    presentation_id: Optional[str] = None
) -> Dict:
    """Add a chart to a slide with comprehensive error handling."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Check if slide index is valid
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate chart type
    valid_chart_types = [
        'column', 'stacked_column', 'bar', 'stacked_bar', 'line', 
        'line_markers', 'pie', 'doughnut', 'area', 'stacked_area', 
        'scatter', 'radar', 'radar_markers'
    ]
    if chart_type.lower() not in valid_chart_types:
        return {
            "error": f"Invalid chart type: '{chart_type}'. Valid types are: {', '.join(valid_chart_types)}"
        }
    
    # Validate series data
    if len(series_names) != len(series_values):
        return {
            "error": f"Number of series names ({len(series_names)}) must match number of series values ({len(series_values)})"
        }
    
    # Validate categories list
    if not categories:
        return {
            "error": "Categories list cannot be empty"
        }
    
    # Validate that all series have the same number of values as categories
    for i, values in enumerate(series_values):
        if len(values) != len(categories):
            return {
                "error": f"Series '{series_names[i]}' has {len(values)} values but there are {len(categories)} categories"
            }
    
    try:
        # Add the chart
        chart, error = ppt_utils.safe_operation(
            "add_chart",
            lambda: ppt_utils.add_chart(
                slide, chart_type, left, top, width, height,
                categories, series_names, series_values
            )
        )
        
        if error:
            return {"error": error}
        
        # Format the chart
        _, error = ppt_utils.safe_operation(
            "format_chart",
            lambda: ppt_utils.format_chart(
                chart,
                has_legend=has_legend,
                legend_position=legend_position,
                has_data_labels=has_data_labels,
                title=title
            )
        )
        
        if error:
            return {
                "warning": f"Chart created but failed to format: {error}",
                "shape_index": len(slide.shapes) - 1
            }
        
        return {
            "message": f"Added {chart_type} chart to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1
        }
    except Exception as e:
        return {
            "error": f"Failed to add chart: {str(e)}"
        }

@app.tool()
def format_text_advanced(
    presentation_id: str = None,
    slide_index: int = 0,
    shape_index: int = 0,
    font_size: int = None,
    font_name: str = None,
    bold: bool = None,
    italic: bool = None,
    color: List[int] = None,
    alignment: str = None,
    auto_resize: bool = True,
    min_font_size: int = 8,
    max_font_size: int = None
) -> Dict[str, Any]:
    """
    Apply advanced text formatting with automatic overflow handling and font size adjustment.
    
    This tool provides enhanced text formatting capabilities including:
    - Automatic font size adjustment when text overflows
    - Text wrapping and container optimization
    - Robust error handling for font modification issues
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        slide_index: Index of the slide (0-based)
        shape_index: Index of the shape containing text (0-based)
        font_size: Font size in points
        font_name: Font name (e.g., 'Arial', 'Times New Roman')
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color as [r, g, b] list (0-255 range)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        auto_resize: Whether to automatically reduce font size if text overflows
        min_font_size: Minimum font size when auto-resizing (default: 8)
        max_font_size: Maximum font size when auto-resizing (default: original font_size)
    
    Returns:
        Dictionary with formatting results including final font size and warnings
    
    Examples:
        Format text with auto-resize enabled:
            format_text_advanced(slide_index=0, shape_index=0, font_size=24, auto_resize=True)
            
        Apply bold formatting with specific color:
            format_text_advanced(slide_index=0, shape_index=0, bold=True, color=[255, 0, 0])
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Presentation has {len(pres.slides)} slides"}
        
        slide = pres.slides[slide_index]
        
        # Validate shape index
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}. Slide has {len(slide.shapes)} shapes"}
        
        shape = slide.shapes[shape_index]
        
        # Check if shape has text
        if not hasattr(shape, 'text_frame'):
            return {"error": f"Shape at index {shape_index} does not contain text"}
        
        # Apply advanced text formatting
        result = ppt_utils.format_text_advanced(
            shape.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=tuple(color) if color else None,
            alignment=alignment,
            auto_resize=auto_resize,
            min_font_size=min_font_size,
            max_font_size=max_font_size
        )
        
        if result['success']:
            message = f"Applied advanced formatting to shape {shape_index} on slide {slide_index}"
            if result['auto_resized']:
                message += f" (font size adjusted from {result['original_font_size']} to {result['final_font_size']})"
            
            return {
                "message": message,
                "formatting_result": result
            }
        else:
            return {"error": f"Failed to format text: {result.get('error', 'Unknown error')}"}
            
    except Exception as e:
        return {"error": f"Failed to format text: {str(e)}"}

@app.tool()
def validate_text_fit(
    presentation_id: str = None,
    slide_index: int = 0,
    shape_index: int = 0,
    text_content: str = None,
    font_size: int = 12
) -> Dict[str, Any]:
    """
    Validate whether text content will fit in a shape container and get optimization suggestions.
    
    This tool helps prevent text overflow issues by:
    - Estimating if text will fit in the current container
    - Suggesting optimal font sizes
    - Recommending container dimension adjustments
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        slide_index: Index of the slide (0-based)
        shape_index: Index of the shape containing text (0-based)
        text_content: Text content to validate (uses current text if not specified)
        font_size: Font size to test (in points)
    
    Returns:
        Dictionary with validation results and optimization suggestions
    
    Examples:
        Check if current text fits:
            validate_text_fit(slide_index=0, shape_index=0)
            
        Test specific text with font size:
            validate_text_fit(slide_index=0, shape_index=0, text_content="Long text here", font_size=16)
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Presentation has {len(pres.slides)} slides"}
        
        slide = pres.slides[slide_index]
        
        # Validate shape index
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}. Slide has {len(slide.shapes)} shapes"}
        
        shape = slide.shapes[shape_index]
        
        # Check if shape has text
        if not hasattr(shape, 'text_frame'):
            return {"error": f"Shape at index {shape_index} does not contain text"}
        
        # Use current text if none specified
        if text_content is None:
            text_content = shape.text_frame.text
        
        # Validate text container
        validation_result = ppt_utils.validate_text_container(shape, text_content, font_size)
        
        return {
            "message": f"Validated text fit for shape {shape_index} on slide {slide_index}",
            "validation_result": validation_result,
            "text_length": len(text_content),
            "tested_font_size": font_size
        }
        
    except Exception as e:
        return {"error": f"Failed to validate text fit: {str(e)}"}

@app.tool()
def validate_and_fix_slide(
    presentation_id: str = None,
    slide_index: int = 0,
    auto_fix: bool = True,
    min_font_size: int = 8,
    max_font_size: int = 72
) -> Dict[str, Any]:
    """
    Comprehensively validate and automatically fix slide content issues.
    
    This tool performs thorough validation of:
    - Font sizes (too large/small for readability and layout)
    - Text boundary overflow and content fitting
    - Shape positioning and layout reasonableness
    - Text margins and spacing
    - Shape overlaps and proper spacing
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        slide_index: Index of the slide to validate (0-based)
        auto_fix: Whether to automatically fix detected issues
        min_font_size: Minimum allowed font size in points
        max_font_size: Maximum allowed font size in points
    
    Returns:
        Dictionary with detailed validation results and applied fixes
    
    Examples:
        Validate and fix current slide:
            validate_and_fix_slide(slide_index=0)
            
        Check slide without applying fixes:
            validate_and_fix_slide(slide_index=1, auto_fix=False)
            
        Custom font size constraints:
            validate_and_fix_slide(slide_index=0, min_font_size=10, max_font_size=48)
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Presentation has {len(pres.slides)} slides"}
        
        slide = pres.slides[slide_index]
        
        # Parameter validation
        valid, error_msg = validate_parameters({
            'min_font_size': (min_font_size, [(is_positive, "Minimum font size must be positive")]),
            'max_font_size': (max_font_size, [(is_positive, "Maximum font size must be positive")]),
            'slide_index': (slide_index, [(is_non_negative, "Slide index must be non-negative")])
        })
        
        if not valid:
            return {"error": error_msg}
        
        if min_font_size >= max_font_size:
            return {"error": "Minimum font size must be less than maximum font size"}
        
        # Perform validation and fixing
        result = ppt_utils.validate_and_fix_slide_content(
            slide, 
            auto_fix=auto_fix,
            min_font_size=min_font_size,
            max_font_size=max_font_size
        )
        
        # Enhance the result with additional context
        result['slide_info'] = {
            'slide_index': slide_index,
            'total_slides': len(pres.slides),
            'auto_fix_enabled': auto_fix,
            'font_size_constraints': {
                'min': min_font_size,
                'max': max_font_size
            }
        }
        
        return result
        
    except Exception as e:
        return {"error": f"Failed to validate slide: {str(e)}"}

@app.tool()
def add_textbox_advanced(
    presentation_id: str = None,
    slide_index: int = 0,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 4.0,
    height: float = 2.0,
    text: str = "Sample text",
    font_size: int = 12,
    font_name: str = None,
    bold: bool = None,
    italic: bool = None,
    color: List[int] = None,
    alignment: str = "left",
    auto_resize: bool = True
) -> Dict[str, Any]:
    """
    Add a textbox with advanced formatting and automatic overflow handling.
    
    This enhanced textbox creation tool provides:
    - Automatic font size adjustment to prevent overflow
    - Built-in text formatting options
    - Intelligent container sizing
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        slide_index: Index of the slide (0-based)
        left: Left position in inches
        top: Top position in inches  
        width: Width in inches
        height: Height in inches
        text: Text content for the textbox
        font_size: Font size in points
        font_name: Font name (e.g., 'Arial', 'Times New Roman')
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color as [r, g, b] list (0-255 range)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        auto_resize: Whether to automatically resize font if text overflows
    
    Returns:
        Dictionary with creation results and shape information
    
    Examples:
        Add a simple textbox with auto-resize:
            add_textbox_advanced(slide_index=0, text="Hello World", font_size=16, auto_resize=True)
            
        Add a formatted textbox with custom styling:
            add_textbox_advanced(slide_index=0, text="Important Note", font_size=14, bold=True, color=[255, 0, 0])
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Presentation has {len(pres.slides)} slides"}
        
        slide = pres.slides[slide_index]
        
        # Validate parameters
        params = {
            'left': (left, [(is_non_negative, "must be non-negative")]),
            'top': (top, [(is_non_negative, "must be non-negative")]),
            'width': (width, [(is_positive, "must be positive")]),
            'height': (height, [(is_positive, "must be positive")]),
            'font_size': (font_size, [(is_positive, "must be positive")])
        }
        
        if color:
            params['color'] = (color, [(is_valid_rgb, "must be [r, g, b] with values 0-255")])
        
        valid, error_msg = validate_parameters(params)
        if not valid:
            return {"error": error_msg}
        
        # Add textbox with advanced formatting
        textbox = ppt_utils.add_textbox(
            slide, left, top, width, height, text,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=tuple(color) if color else None,
            alignment=alignment,
            auto_resize=auto_resize
        )
        
        return {
            "message": f"Added advanced textbox to slide {slide_index}",
            "shape_index": len(slide.shapes) - 1,
            "dimensions": {
                "left": left,
                "top": top,
                "width": width,
                "height": height
            },
            "formatting": {
                "font_size": font_size,
                "font_name": font_name,
                "bold": bold,
                "italic": italic,
                "color": color,
                "alignment": alignment,
                "auto_resize": auto_resize
            }
        }
        
    except Exception as e:
        return {"error": f"Failed to add advanced textbox: {str(e)}"}

@app.tool()
def add_professional_slide(
    presentation_id: str = None,
    slide_type: str = 'title_content',
    color_scheme: str = 'modern_blue',
    title: str = None,
    content: List[str] = None
) -> Dict[str, Any]:
    """
    Add a professionally designed slide with modern styling and proper layout.
    
    This tool creates slides with:
    - Professional color schemes and typography
    - Proper spacing and margins
    - Modern flat design principles
    - Consistent visual hierarchy
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        slide_type: Type of slide ('title', 'title_content', 'content', 'two_column', 'blank')
        color_scheme: Color scheme ('modern_blue', 'corporate_gray', 'elegant_green', 'warm_red')
        title: Optional title text to add immediately
        content: Optional list of bullet points or content to add
    
    Returns:
        Dictionary with slide creation results and design information
    
    Examples:
        Create a professional title slide:
            add_professional_slide(slide_type='title', title='Quarterly Results', color_scheme='corporate_gray')
            
        Create a content slide with bullet points:
            add_professional_slide(slide_type='title_content', title='Key Points', 
                                  content=['Point 1', 'Point 2'], color_scheme='modern_blue')
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Add professional slide
        slide, layout, design_info = ppt_utils.add_professional_slide(
            pres, slide_type, color_scheme
        )
        
        # Add title if provided
        if title:
            ppt_utils.set_professional_title(slide, title, color_scheme)
        
        # Add content if provided
        if content and slide_type in ['title_content', 'content']:
            try:
                # Find content placeholder (usually index 1)
                content_placeholder = None
                for i, placeholder in enumerate(slide.placeholders):
                    if i == 1:  # Content placeholder
                        content_placeholder = placeholder
                        break
                
                if content_placeholder:
                    ppt_utils.add_professional_bullet_points(
                        content_placeholder, content, color_scheme, hierarchical=True
                    )
            except Exception as e:
                design_info['content_warning'] = f"Could not add content: {str(e)}"
        
        return {
            "message": f"Added professional {slide_type} slide",
            "slide_index": len(pres.slides) - 1,
            "design_info": design_info
        }
        
    except Exception as e:
        return {"error": f"Failed to add professional slide: {str(e)}"}

@app.tool()
def apply_professional_theme(
    presentation_id: str = None,
    color_scheme: str = 'modern_blue',
    apply_to_existing: bool = True
) -> Dict[str, Any]:
    """
    Apply a professional theme and color scheme to the presentation.
    
    This tool enhances presentation quality by:
    - Applying consistent professional color schemes
    - Setting up proper typography standards
    - Configuring modern design principles
    - Optionally updating existing slides
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        color_scheme: Color scheme to apply ('modern_blue', 'corporate_gray', 'elegant_green', 'warm_red')
        apply_to_existing: Whether to update formatting of existing slides
    
    Returns:
        Dictionary with theming results and applied settings
    
    Examples:
        Apply modern blue theme:
            apply_professional_theme(color_scheme='modern_blue')
            
        Apply corporate theme without changing existing slides:
            apply_professional_theme(color_scheme='corporate_gray', apply_to_existing=False)
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Get color scheme info
        if color_scheme not in ppt_utils.PROFESSIONAL_COLOR_SCHEMES:
            return {"error": f"Unknown color scheme: {color_scheme}. Available: {list(ppt_utils.PROFESSIONAL_COLOR_SCHEMES.keys())}"}
        
        scheme_colors = ppt_utils.PROFESSIONAL_COLOR_SCHEMES[color_scheme]
        updated_slides = 0
        warnings = []
        
        # Apply theme to existing slides if requested
        if apply_to_existing:
            for i, slide in enumerate(pres.slides):
                try:
                    # Apply background
                    ppt_utils.apply_professional_slide_background(slide, color_scheme)
                    
                    # Update title formatting if present
                    if slide.shapes.title and slide.shapes.title.text:
                        ppt_utils.set_professional_title(slide, slide.shapes.title.text, color_scheme)
                    
                    # Update text formatting in shapes
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.text_frame.text:
                            try:
                                ppt_utils.format_text_advanced(
                                    shape.text_frame,
                                    color=ppt_utils.get_professional_color(color_scheme, 'text'),
                                    auto_resize=True
                                )
                            except:
                                continue
                    
                    updated_slides += 1
                    
                except Exception as e:
                    warnings.append(f"Slide {i}: {str(e)}")
                    continue
        
        return {
            "message": f"Applied {color_scheme} theme to presentation",
            "color_scheme": color_scheme,
            "scheme_colors": scheme_colors,
            "updated_slides": updated_slides,
            "total_slides": len(pres.slides),
            "warnings": warnings if warnings else None
        }
        
    except Exception as e:
        return {"error": f"Failed to apply professional theme: {str(e)}"}

@app.tool()
def enhance_existing_slide(
    presentation_id: str = None,
    slide_index: int = 0,
    color_scheme: str = 'modern_blue',
    enhance_title: bool = True,
    enhance_content: bool = True,
    enhance_shapes: bool = True,
    enhance_charts: bool = True
) -> Dict[str, Any]:
    """
    Enhance an existing slide with professional formatting and modern design.
    
    This tool improves slide quality by:
    - Applying professional typography and colors
    - Modernizing shape and chart styling
    - Optimizing spacing and layout
    - Ensuring visual consistency
    
    Args:
        presentation_id: ID of the presentation (uses current if not specified)
        slide_index: Index of the slide to enhance (0-based)
        color_scheme: Color scheme to apply
        enhance_title: Whether to enhance title formatting
        enhance_content: Whether to enhance text content formatting
        enhance_shapes: Whether to enhance shape formatting
        enhance_charts: Whether to enhance chart formatting
    
    Returns:
        Dictionary with enhancement results and changes made
    
    Examples:
        Enhance all elements of a slide:
            enhance_existing_slide(slide_index=0, color_scheme='corporate_gray')
            
        Only enhance text formatting:
            enhance_existing_slide(slide_index=1, enhance_shapes=False, enhance_charts=False)
    """
    try:
        # Get the presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Presentation has {len(pres.slides)} slides"}
        
        slide = pres.slides[slide_index]
        enhancements = []
        warnings = []
        
        # Enhance slide background
        try:
            ppt_utils.apply_professional_slide_background(slide, color_scheme)
            enhancements.append("Applied professional background")
        except Exception as e:
            warnings.append(f"Background: {str(e)}")
        
        # Enhance title
        if enhance_title and slide.shapes.title and slide.shapes.title.text:
            try:
                original_title = slide.shapes.title.text
                ppt_utils.set_professional_title(slide, original_title, color_scheme)
                enhancements.append("Enhanced title formatting")
            except Exception as e:
                warnings.append(f"Title: {str(e)}")
        
        # Enhance content and shapes
        shape_count = 0
        chart_count = 0
        
        for i, shape in enumerate(slide.shapes):
            try:
                # Skip title shape (already handled)
                if shape == slide.shapes.title:
                    continue
                
                # Enhance text content
                if enhance_content and hasattr(shape, 'text_frame') and shape.text_frame.text:
                    try:
                        ppt_utils.format_text_advanced(
                            shape.text_frame,
                            font_name=ppt_utils.get_professional_font('body', 'medium')['name'],
                            color=ppt_utils.get_professional_color(color_scheme, 'text'),
                            auto_resize=True
                        )
                        enhancements.append(f"Enhanced text in shape {i}")
                    except Exception as e:
                        warnings.append(f"Text shape {i}: {str(e)}")
                
                # Enhance regular shapes
                if enhance_shapes and hasattr(shape, 'fill') and not hasattr(shape, 'chart'):
                    try:
                        ppt_utils.format_professional_shape(shape, color_scheme, 'primary')
                        shape_count += 1
                    except Exception as e:
                        warnings.append(f"Shape {i}: {str(e)}")
                
                # Enhance charts
                if enhance_charts and hasattr(shape, 'chart'):
                    try:
                        ppt_utils.format_professional_chart(shape.chart, color_scheme)
                        chart_count += 1
                    except Exception as e:
                        warnings.append(f"Chart {i}: {str(e)}")
                        
            except Exception as e:
                warnings.append(f"Shape {i}: {str(e)}")
                continue
        
        if shape_count > 0:
            enhancements.append(f"Enhanced {shape_count} shapes")
        if chart_count > 0:
            enhancements.append(f"Enhanced {chart_count} charts")
        
        return {
            "message": f"Enhanced slide {slide_index} with {color_scheme} theme",
            "slide_index": slide_index,
            "color_scheme": color_scheme,
            "enhancements": enhancements,
            "shapes_enhanced": shape_count,
            "charts_enhanced": chart_count,
            "warnings": warnings if warnings else None
        }
        
    except Exception as e:
        return {"error": f"Failed to enhance slide: {str(e)}"}

@app.tool()
def get_color_schemes() -> Dict[str, Any]:
    """
    Get all available professional color schemes with their color values.
    
    Returns:
        Dictionary containing all available color schemes and their color definitions
    
    Examples:
        View available color schemes:
            get_color_schemes()
    """
    return {
        "message": "Available professional color schemes",
        "color_schemes": ppt_utils.PROFESSIONAL_COLOR_SCHEMES,
        "font_settings": ppt_utils.PROFESSIONAL_FONTS,
        "layout_constants": ppt_utils.PROFESSIONAL_LAYOUT
    }

# ---- Advanced Features: Gradient Backgrounds, Image Enhancement, Font Beautification ----

@app.tool()
def set_slide_gradient_background(slide_index: int, start_color: List[int], end_color: List[int], 
                                 direction: str = "horizontal", presentation_id: Optional[str] = None) -> Dict[str, Any]:
    """
    Set a gradient background for a specific slide using Pillow to generate the gradient.
    
    Args:
        slide_index: Index of the slide (0-based)
        start_color: RGB color for gradient start [r, g, b] (0-255 each)
        end_color: RGB color for gradient end [r, g, b] (0-255 each)
        direction: Gradient direction ('horizontal', 'vertical', 'diagonal')
        presentation_id: ID of the presentation (uses current if not specified)
        
    Returns:
        Dictionary with operation results
        
    Examples:
        Set horizontal blue to white gradient on slide 0:
            set_slide_gradient_background(0, [0, 120, 215], [255, 255, 255], "horizontal")
            
        Set diagonal gradient from red to yellow on slide 1:
            set_slide_gradient_background(1, [220, 20, 60], [255, 215, 0], "diagonal")
    """
    try:
        # Parameter validation
        valid, error = validate_parameters({
            'slide_index': (slide_index, [(is_non_negative, "must be non-negative")]),
            'start_color': (start_color, [(is_valid_rgb, "must be a list of 3 integers (0-255)")]),
            'end_color': (end_color, [(is_valid_rgb, "must be a list of 3 integers (0-255)")]),
            'direction': (direction, [(is_in_list(['horizontal', 'vertical', 'diagonal']), 
                                    "must be 'horizontal', 'vertical', or 'diagonal'")])
        })
        
        if not valid:
            return {"error": error}
        
        # Get presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation with ID '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index >= len(pres.slides):
            return {"error": f"Slide index {slide_index} is out of range. Presentation has {len(pres.slides)} slides."}
        
        slide = pres.slides[slide_index]
        
        # Convert RGB lists to tuples
        start_color_tuple = tuple(start_color)
        end_color_tuple = tuple(end_color)
        
        # Set gradient background
        result = ppt_utils.set_slide_gradient_background(slide, start_color_tuple, end_color_tuple, direction)
        
        if result['success']:
            return {
                "message": result['message'],
                "slide_index": slide_index,
                "gradient_info": result['gradient_info'],
                "presentation_id": presentation_id or current_presentation_id
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to set gradient background: {str(e)}"}

@app.tool()
def create_professional_gradient_background(slide_index: int, color_scheme: str = "modern_blue", 
                                          style: str = "subtle", direction: str = "diagonal",
                                          presentation_id: Optional[str] = None) -> Dict[str, Any]:
    """
    Create a professional gradient background using predefined color schemes.
    
    Args:
        slide_index: Index of the slide (0-based)
        color_scheme: Professional color scheme ('modern_blue', 'corporate_gray', 'elegant_green', 'warm_red')
        style: Gradient style ('subtle', 'bold', 'accent')
        direction: Gradient direction ('horizontal', 'vertical', 'diagonal')
        presentation_id: ID of the presentation (uses current if not specified)
        
    Returns:
        Dictionary with operation results
        
    Examples:
        Create subtle blue gradient background:
            create_professional_gradient_background(0, "modern_blue", "subtle", "diagonal")
            
        Create bold corporate gradient:
            create_professional_gradient_background(1, "corporate_gray", "bold", "horizontal")
    """
    try:
        # Parameter validation
        valid_schemes = list(ppt_utils.PROFESSIONAL_COLOR_SCHEMES.keys())
        valid_styles = ['subtle', 'bold', 'accent']
        valid_directions = ['horizontal', 'vertical', 'diagonal']
        
        valid, error = validate_parameters({
            'slide_index': (slide_index, [(is_non_negative, "must be non-negative")]),
            'color_scheme': (color_scheme, [(is_in_list(valid_schemes), f"must be one of {valid_schemes}")]),
            'style': (style, [(is_in_list(valid_styles), f"must be one of {valid_styles}")]),
            'direction': (direction, [(is_in_list(valid_directions), f"must be one of {valid_directions}")])
        })
        
        if not valid:
            return {"error": error}
        
        # Get presentation
        if presentation_id:
            if presentation_id not in presentations:
                return {"error": f"Presentation with ID '{presentation_id}' not found"}
            pres = presentations[presentation_id]
        else:
            pres = get_current_presentation()
        
        # Validate slide index
        if slide_index >= len(pres.slides):
            return {"error": f"Slide index {slide_index} is out of range. Presentation has {len(pres.slides)} slides."}
        
        slide = pres.slides[slide_index]
        
        # Create professional gradient background
        result = ppt_utils.create_professional_gradient_background(slide, color_scheme, style, direction)
        
        if result['success']:
            return {
                "message": result['message'],
                "slide_index": slide_index,
                "color_scheme": color_scheme,
                "style": style,
                "gradient_info": result['gradient_info'],
                "presentation_id": presentation_id or current_presentation_id
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to create professional gradient background: {str(e)}"}

@app.tool()
def enhance_image_with_pillow(image_path: str, brightness: float = 1.0, contrast: float = 1.0,
                             saturation: float = 1.0, sharpness: float = 1.0, blur_radius: float = 0,
                             filter_type: Optional[str] = None, output_path: Optional[str] = None) -> Dict[str, Any]:
    """
    Enhance an image using Pillow filters and adjustments.
    
    Args:
        image_path: Path to the input image file
        brightness: Brightness adjustment (1.0 = no change, >1.0 = brighter, <1.0 = darker)
        contrast: Contrast adjustment (1.0 = no change, >1.0 = more contrast)
        saturation: Color saturation adjustment (1.0 = no change, >1.0 = more saturated)
        sharpness: Sharpness adjustment (1.0 = no change, >1.0 = sharper)
        blur_radius: Gaussian blur radius (0 = no blur, >0 = blur effect)
        filter_type: Special filter to apply ('DETAIL', 'EDGE_ENHANCE', 'EMBOSS', 'SMOOTH', 'SHARPEN', etc.)
        output_path: Path for enhanced image (optional, creates temp file if not provided)
        
    Returns:
        Dictionary with enhancement results and output path
        
    Examples:
        Enhance image for presentation (brighter, more contrast):
            enhance_image_with_pillow("/path/to/image.jpg", brightness=1.1, contrast=1.2, sharpness=1.1)
            
        Apply artistic effect:
            enhance_image_with_pillow("/path/to/image.jpg", saturation=1.5, filter_type="EDGE_ENHANCE")
    """
    try:
        # Parameter validation
        valid_filters = ['DETAIL', 'EDGE_ENHANCE', 'EDGE_ENHANCE_MORE', 'EMBOSS', 'FIND_EDGES', 
                        'SMOOTH', 'SMOOTH_MORE', 'SHARPEN']
        
        valid, error = validate_parameters({
            'brightness': (brightness, [(lambda x: x > 0, "must be positive")]),
            'contrast': (contrast, [(lambda x: x > 0, "must be positive")]),
            'saturation': (saturation, [(lambda x: x > 0, "must be positive")]),
            'sharpness': (sharpness, [(lambda x: x > 0, "must be positive")]),
            'blur_radius': (blur_radius, [(is_non_negative, "must be non-negative")])
        })
        
        if not valid:
            return {"error": error}
        
        if filter_type and filter_type not in valid_filters:
            return {"error": f"Invalid filter_type. Must be one of: {valid_filters}"}
        
        # Check if image file exists
        if not os.path.exists(image_path):
            return {"error": f"Image file not found: {image_path}"}
        
        # Prepare enhancement settings
        enhancements = {
            'brightness': brightness,
            'contrast': contrast,
            'saturation': saturation,
            'sharpness': sharpness,
            'blur_radius': blur_radius,
            'filter': filter_type
        }
        
        # Enhance the image
        result = ppt_utils.enhance_image_pillow(image_path, enhancements, output_path)
        
        if result['success']:
            return {
                "message": result['message'],
                "original_path": result['original_path'],
                "enhanced_path": result['enhanced_path'],
                "enhancements_applied": result['enhancements_applied'],
                "image_info": result.get('image_info', {})
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to enhance image: {str(e)}"}

@app.tool()
def apply_professional_image_enhancement(image_path: str, style: str = "presentation", 
                                       output_path: Optional[str] = None) -> Dict[str, Any]:
    """
    Apply professional image enhancement presets suitable for presentations.
    
    Args:
        image_path: Path to the input image file
        style: Enhancement style ('presentation', 'vibrant', 'subtle', 'sharp')
        output_path: Path for enhanced image (optional, creates temp file if not provided)
        
    Returns:
        Dictionary with enhancement results
        
    Examples:
        Enhance image for business presentation:
            apply_professional_image_enhancement("/path/to/image.jpg", "presentation")
            
        Create vibrant image for creative presentation:
            apply_professional_image_enhancement("/path/to/image.jpg", "vibrant")
    """
    try:
        # Parameter validation
        valid_styles = ['presentation', 'vibrant', 'subtle', 'sharp']
        
        valid, error = validate_parameters({
            'style': (style, [(is_in_list(valid_styles), f"must be one of {valid_styles}")])
        })
        
        if not valid:
            return {"error": error}
        
        # Check if image file exists
        if not os.path.exists(image_path):
            return {"error": f"Image file not found: {image_path}"}
        
        # Apply professional enhancement
        result = ppt_utils.apply_professional_image_enhancement(image_path, style)
        
        if result['success']:
            return {
                "message": result['message'],
                "original_path": result['original_path'],
                "enhanced_path": result['enhanced_path'],
                "style": style,
                "enhancements_applied": result['enhancements_applied'],
                "image_info": result.get('image_info', {})
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply professional image enhancement: {str(e)}"}

@app.tool()
def analyze_font_file(font_path: str) -> Dict[str, Any]:
    """
    Analyze a font file using FontTools to extract detailed information.
    
    Args:
        font_path: Path to the font file (.ttf, .otf)
        
    Returns:
        Dictionary with comprehensive font information
        
    Examples:
        Analyze a font file:
            analyze_font_file("/path/to/font.ttf")
    """
    try:
        # Check if font file exists
        if not os.path.exists(font_path):
            return {"error": f"Font file not found: {font_path}"}
        
        # Analyze the font
        result = ppt_utils.analyze_font_file(font_path)
        
        if result['success']:
            return {
                "message": result['message'],
                "font_info": result['font_info']
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to analyze font: {str(e)}"}

@app.tool()
def optimize_font_for_presentation(font_path: str, output_path: Optional[str] = None, 
                                  text_content: Optional[str] = None) -> Dict[str, Any]:
    """
    Optimize a font file for presentation use by subsetting and optimizing.
    
    Args:
        font_path: Path to the input font file (.ttf, .otf)
        output_path: Path for optimized font (optional, creates temp file if not provided)
        text_content: Specific text content to optimize for (optional, uses common characters if not provided)
        
    Returns:
        Dictionary with optimization results
        
    Examples:
        Optimize font for general presentation use:
            optimize_font_for_presentation("/path/to/font.ttf")
            
        Optimize font for specific content:
            optimize_font_for_presentation("/path/to/font.ttf", text_content="Hello World Presentation")
    """
    try:
        # Check if font file exists
        if not os.path.exists(font_path):
            return {"error": f"Font file not found: {font_path}"}
        
        # Optimize the font
        result = ppt_utils.optimize_font_for_presentation(font_path, output_path, text_content)
        
        if result['success']:
            return {
                "message": result['message'],
                "original_path": result['original_path'],
                "optimized_path": result['optimized_path'],
                "size_reduction_percent": result['size_reduction'],
                "file_sizes": result['file_sizes']
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to optimize font: {str(e)}"}

@app.tool()
def get_font_recommendations(font_path: str, presentation_type: str = "business") -> Dict[str, Any]:
    """
    Get font usage recommendations based on analysis and presentation type.
    
    Args:
        font_path: Path to the font file to analyze
        presentation_type: Type of presentation ('business', 'creative', 'academic')
        
    Returns:
        Dictionary with font recommendations and usage suggestions
        
    Examples:
        Get recommendations for business presentation:
            get_font_recommendations("/path/to/font.ttf", "business")
            
        Get recommendations for creative presentation:
            get_font_recommendations("/path/to/font.ttf", "creative")
    """
    try:
        # Parameter validation
        valid_types = ['business', 'creative', 'academic']
        
        valid, error = validate_parameters({
            'presentation_type': (presentation_type, [(is_in_list(valid_types), f"must be one of {valid_types}")])
        })
        
        if not valid:
            return {"error": error}
        
        # Check if font file exists
        if not os.path.exists(font_path):
            return {"error": f"Font file not found: {font_path}"}
        
        # First analyze the font
        analysis_result = ppt_utils.analyze_font_file(font_path)
        
        if not analysis_result['success']:
            return {"error": f"Failed to analyze font: {analysis_result['message']}"}
        
        # Get recommendations based on analysis
        recommendations = ppt_utils.get_font_recommendations(analysis_result, presentation_type)
        
        return {
            "message": f"Font recommendations for {presentation_type} presentation",
            "font_path": font_path,
            "presentation_type": presentation_type,
            "font_analysis": analysis_result['font_info'],
            "recommendations": recommendations
        }
            
    except Exception as e:
        return {"error": f"Failed to get font recommendations: {str(e)}"}

# ---- Picture Effects Tools ----

@app.tool()
def apply_picture_shadow(
    slide_index: int,
    shape_index: int,
    shadow_type: str = 'outer',
    blur_radius: float = 4.0,
    distance: float = 3.0,
    direction: float = 315.0,
    color: List[int] = [0, 0, 0],
    transparency: float = 0.6,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply shadow effect to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid, error = validate_parameters({
        'blur_radius': (blur_radius, [(is_non_negative, "must be non-negative")]),
        'distance': (distance, [(is_non_negative, "must be non-negative")]),
        'direction': (direction, [(is_in_range(0, 360), "must be between 0 and 360 degrees")]),
        'transparency': (transparency, [(is_in_range(0, 1), "must be between 0.0 and 1.0")]),
        'color': (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_shadow(
            shape, shadow_type, blur_radius, distance, direction, tuple(color), transparency
        )
        
        if result['success']:
            return {
                "message": result['message'],
                "shadow_info": result['shadow_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply shadow effect: {str(e)}"}

@app.tool()
def apply_picture_reflection(
    slide_index: int,
    shape_index: int,
    size: float = 0.5,
    transparency: float = 0.5,
    distance: float = 0.0,
    blur: float = 4.0,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply reflection effect to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid, error = validate_parameters({
        'size': (size, [(is_in_range(0, 1), "must be between 0.0 and 1.0")]),
        'transparency': (transparency, [(is_in_range(0, 1), "must be between 0.0 and 1.0")]),
        'distance': (distance, [(is_non_negative, "must be non-negative")]),
        'blur': (blur, [(is_non_negative, "must be non-negative")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_reflection(shape, size, transparency, distance, blur)
        
        if result['success']:
            return {
                "message": result['message'],
                "reflection_info": result['reflection_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply reflection effect: {str(e)}"}

@app.tool()
def apply_picture_glow(
    slide_index: int,
    shape_index: int,
    size: float = 5.0,
    color: List[int] = [0, 176, 240],
    transparency: float = 0.4,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply glow effect to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid, error = validate_parameters({
        'size': (size, [(is_positive, "must be positive")]),
        'transparency': (transparency, [(is_in_range(0, 1), "must be between 0.0 and 1.0")]),
        'color': (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_glow(shape, size, tuple(color), transparency)
        
        if result['success']:
            return {
                "message": result['message'],
                "glow_info": result['glow_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply glow effect: {str(e)}"}

@app.tool()
def apply_picture_soft_edges(
    slide_index: int,
    shape_index: int,
    radius: float = 2.5,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply soft edges effect to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid, error = validate_parameters({
        'radius': (radius, [(is_non_negative, "must be non-negative")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_soft_edges(shape, radius)
        
        if result['success']:
            return {
                "message": result['message'],
                "soft_edge_info": result['soft_edge_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply soft edges effect: {str(e)}"}

@app.tool()
def apply_picture_rotation(
    slide_index: int,
    shape_index: int,
    rotation: float,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply rotation to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    try:
        result = ppt_utils.apply_picture_rotation(shape, rotation)
        
        if result['success']:
            return {
                "message": result['message'],
                "rotation_info": result['rotation_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply rotation: {str(e)}"}

@app.tool()
def apply_picture_transparency(
    slide_index: int,
    shape_index: int,
    transparency: float,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply transparency to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid, error = validate_parameters({
        'transparency': (transparency, [(is_in_range(0, 1), "must be between 0.0 and 1.0")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_transparency(shape, transparency)
        
        if result['success']:
            return {
                "message": result['message'],
                "transparency_info": result['transparency_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply transparency: {str(e)}"}

@app.tool()
def apply_picture_bevel(
    slide_index: int,
    shape_index: int,
    bevel_type: str = 'circle',
    width: float = 6.0,
    height: float = 6.0,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply bevel effect to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid_bevel_types = ['circle', 'square', 'slope', 'riblet']
    valid, error = validate_parameters({
        'bevel_type': (bevel_type, [(is_in_list(valid_bevel_types), f"must be one of {valid_bevel_types}")]),
        'width': (width, [(is_positive, "must be positive")]),
        'height': (height, [(is_positive, "must be positive")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_bevel(shape, bevel_type, width, height)
        
        if result['success']:
            return {
                "message": result['message'],
                "bevel_info": result['bevel_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply bevel effect: {str(e)}"}

@app.tool()
def apply_picture_filter(
    slide_index: int,
    shape_index: int,
    filter_type: str = 'none',
    intensity: float = 0.5,
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply color filter effect to a picture shape."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    # Validate parameters
    valid_filter_types = ['grayscale', 'sepia', 'washout', 'none']
    valid, error = validate_parameters({
        'filter_type': (filter_type, [(is_in_list(valid_filter_types), f"must be one of {valid_filter_types}")]),
        'intensity': (intensity, [(is_in_range(0, 1), "must be between 0.0 and 1.0")])
    })
    
    if not valid:
        return {"error": error}
    
    try:
        result = ppt_utils.apply_picture_filter(shape, filter_type, intensity)
        
        if result['success']:
            return {
                "message": result['message'],
                "filter_info": result['filter_info'],
                "slide_index": slide_index,
                "shape_index": shape_index
            }
        else:
            return {"error": result['message']}
            
    except Exception as e:
        return {"error": f"Failed to apply filter: {str(e)}"}

@app.tool()
def apply_combined_picture_effects(
    slide_index: int,
    shape_index: int,
    effects: Dict[str, Dict],
    presentation_id: Optional[str] = None
) -> Dict:
    """Apply multiple picture effects in combination."""
    # Use the specified presentation or the current one
    pres_id = presentation_id if presentation_id is not None else current_presentation_id
    
    if pres_id is None or pres_id not in presentations:
        return {
            "error": "No presentation is currently loaded or the specified ID is invalid"
        }
    
    pres = presentations[pres_id]
    
    # Validate slide index
    if slide_index < 0 or slide_index >= len(pres.slides):
        return {
            "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
        }
    
    slide = pres.slides[slide_index]
    
    # Validate shape index
    if shape_index < 0 or shape_index >= len(slide.shapes):
        return {
            "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
        }
    
    shape = slide.shapes[shape_index]
    
    # Check if shape is a picture
    if not hasattr(shape, 'image') and shape.shape_type.name != 'PICTURE':
        return {
            "error": f"Shape at index {shape_index} is not a picture shape"
        }
    
    try:
        result = ppt_utils.apply_combined_picture_effects(shape, effects)
        
        result['slide_index'] = slide_index
        result['shape_index'] = shape_index
        
        return result
            
    except Exception as e:
        return {"error": f"Failed to apply combined effects: {str(e)}"}

# ---- Main Execution ----
def main(transport: str = "stdio", port: int = 8000):
    if transport == "http":
        import asyncio
        # Set the port for HTTP transport
        app.settings.port = port
        # Start the FastMCP server with HTTP transport
        try:
            app.run(transport='streamable-http')
        except asyncio.exceptions.CancelledError:
            print("Server stopped by user.")
        except KeyboardInterrupt:
            print("Server stopped by user.")
        except Exception as e:
            print(f"Error starting server: {e}")

    else:
        # Run the FastMCP server
        app.run(transport='stdio')

if __name__ == "__main__":
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="MCP Server for PowerPoint manipulation using python-pptx")

    parser.add_argument(
        "-t",
        "--transport",
        type=str,
        default="stdio",
        choices=["stdio", "http"],
        help="Transport method for the MCP server (default: stdio)"
    )

    parser.add_argument(
        "-p",
        "--port",
        type=int,
        default=8000,
        help="Port to run the MCP server on (default: 8000)"
    )
    args = parser.parse_args()
    main(args.transport, args.port)